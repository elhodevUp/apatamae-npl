import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml import parse_xml

def create_expert_presentation():
    prs = Presentation()
    # Format 16:9 Widescreen HD
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # --- PALETTE CHROMATIQUE EXPERTE ---
    BLEU_NUIT       = RGBColor(10, 34, 64)    # #0A2240 (Bleu Institutionnel profond)
    VERT_EMERAUDE   = RGBColor(18, 102, 53)   # #126635 (Vert Territoire dynamique)
    VERT_FORET      = RGBColor(1, 45, 29)     # #012D1D (Vert Sombre Institutionnel)
    OR_SOLEIL       = RGBColor(222, 154, 25)  # #DE9A19 (Or Rayonnement)
    OR_LUMINEUX     = RGBColor(245, 190, 44)  # #F5BE2C (Or Éclatant)
    BLANC_PUR       = RGBColor(255, 255, 255) # #FFFFFF
    FOND_PERLE      = RGBColor(250, 251, 253) # #FAFBFD
    NOIR_CHARBON    = RGBColor(20, 24, 28)    # #14181C
    GRIS_TEXTE      = RGBColor(75, 85, 90)    # #4B555A
    BORDER_SUBTILE  = RGBColor(226, 232, 230) # #E2E8E6
    VERT_DOUX       = RGBColor(235, 246, 239) # #EBF6EF
    BLEU_DOUX       = RGBColor(238, 244, 252) # #EEF4FC
    OR_DOUX         = RGBColor(254, 249, 238) # #FEF9EE

    base_dir = os.path.dirname(os.path.abspath(__file__))
    artifacts_dir = r"C:\Users\LENOVO E15\.gemini\antigravity-ide\brain\2584ad8b-cb21-4a7f-a2cd-17da14691d98"
    e_drive_dir = r"E:\noppal\apatamae"

    # Logo
    logo_path = os.path.join(base_dir, "public", "logo.png")
    if not os.path.exists(logo_path):
        logo_path = os.path.join(e_drive_dir, "logo.png.jpeg")

    # Ultra-Attractive New Images
    img_baobab = os.path.join(artifacts_dir, "hero_cinematic_baobab_1785937820177.png")
    img_women = os.path.join(artifacts_dir, "women_empowerment_vibrant_1785937836631.png")
    img_stem = os.path.join(artifacts_dir, "youth_stem_innovation_1785937852615.png")
    img_agri = os.path.join(artifacts_dir, "sustainable_solar_agriculture_1785937869292.png")
    img_sprout = os.path.join(artifacts_dir, "closing_hands_sprout_glow_1785937883057.png")
    img_dialogue = os.path.join(artifacts_dir, "pptx_intergenerationnel_1785866630658.png")
    img_seminaire = os.path.join(artifacts_dir, "pptx_salles_seminaire_1785866596542.png")
    img_chambre = os.path.join(artifacts_dir, "pptx_chambre_hotes_1785866613303.png")
    img_pres = os.path.join(base_dir, "public", "images", "professional_corporate_portrait_of_a_west_african_woman_leader_elegant_and_screen.png")

    def apply_transition(slide, trans_type="fade", direction=None):
        """Ajoute des transitions XML animées avancées conformes au standard OpenXML"""
        try:
            if direction:
                xml = f'<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:{trans_type} dir="{direction}"/></p:transition>'
            else:
                xml = f'<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:{trans_type}/></p:transition>'
            slide._element.append(parse_xml(xml))
        except Exception:
            pass

    def draw_bg(slide, color=FOND_PERLE):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = color
        bg.line.fill.background()
        return bg

    def add_expert_header(slide, title_text, category_badge="APATAM@E • STRATÉGIE 2026", dark_mode=False):
        """Header professionnel avec badge catégorie, grand titre et ligne or d'accentuation"""
        # Logo miniature
        if os.path.exists(logo_path):
            slide.shapes.add_picture(logo_path, Inches(0.8), Inches(0.35), height=Inches(0.85))

        # Badge catégorie
        tb_cat = slide.shapes.add_textbox(Inches(2.0), Inches(0.35), Inches(8.5), Inches(0.3))
        p_c = tb_cat.text_frame.paragraphs[0]
        p_c.text = category_badge.upper()
        p_c.font.size = Pt(10)
        p_c.font.bold = True
        p_c.font.color.rgb = OR_LUMINEUX if dark_mode else OR_SOLEIL

        # Slogan institutionnel
        p_s = tb_cat.text_frame.add_paragraph()
        p_s.text = "Former, accompagner et innover pour des territoires durables."
        p_s.font.size = Pt(9.5)
        p_s.font.italic = True
        p_s.font.color.rgb = BLANC_PUR if dark_mode else GRIS_TEXTE

        # Grand Titre (Grande écriture 28pt)
        tb_t = slide.shapes.add_textbox(Inches(0.8), Inches(1.25), Inches(11.7), Inches(0.7))
        p_t = tb_t.text_frame.paragraphs[0]
        p_t.text = title_text
        p_t.font.size = Pt(28)
        p_t.font.bold = True
        p_t.font.color.rgb = BLANC_PUR if dark_mode else VERT_FORET

        # Ligne Dorée Subtile
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.95), Inches(11.73), Inches(0.035))
        line.fill.solid()
        line.fill.fore_color.rgb = OR_SOLEIL
        line.line.fill.background()

    def add_signature_bottom_bar(slide, dark_mode=False):
        """Barre inférieure élégante avec les 4 piliers d'action et logo"""
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.65), Inches(11.73), Inches(0.55))
        bar.fill.solid()
        bar.fill.fore_color.rgb = BLEU_NUIT if not dark_mode else VERT_EMERAUDE
        bar.line.fill.background()
        
        btf = bar.text_frame
        btf.word_wrap = True
        p = btf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = "✦  Inclusion Sociale   •   Innovation & Savoirs   •   Gouvernance Territoriale   •   Durabilité Écologique  ✦"
        p.font.size = Pt(10.5)
        p.font.bold = True
        p.font.color.rgb = BLANC_PUR

    # =========================================================================
    # SLIDE 1 : COUVERTURE IMMERSIVE & IMPACTANTE (Transition Zoom)
    # =========================================================================
    s1 = prs.slides.add_slide(blank_layout)
    draw_bg(s1, VERT_FORET)
    apply_transition(s1, "zoom")

    # Image Droite : Baobab Majestueux Haute Résolution sous lumière dorée
    if os.path.exists(img_baobab):
        s1.shapes.add_picture(img_baobab, Inches(6.3), Inches(0.6), width=Inches(6.4))

    # Bloc Gauche
    if os.path.exists(logo_path):
        s1.shapes.add_picture(logo_path, Inches(0.8), Inches(0.8), width=Inches(1.8))

    # Badge Titre
    b_tag = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.7), Inches(4.5), Inches(0.42))
    b_tag.fill.solid()
    b_tag.fill.fore_color.rgb = BLEU_NUIT
    b_tag.line.color.rgb = OR_SOLEIL
    b_tag.line.width = Pt(1)
    ptg = b_tag.text_frame.paragraphs[0]
    ptg.text = "DOCUMENT INSTITUTIONNEL OFFICIEL"
    ptg.font.size = Pt(9.5)
    ptg.font.bold = True
    ptg.font.color.rgb = OR_LUMINEUX
    ptg.alignment = PP_ALIGN.CENTER

    # Grand Titre (52pt)
    tb1 = s1.shapes.add_textbox(Inches(0.75), Inches(3.25), Inches(5.3), Inches(2.8))
    tf1 = tb1.text_frame
    tf1.word_wrap = True

    p1_1 = tf1.paragraphs[0]
    p1_1.text = "APATAM@E"
    p1_1.font.size = Pt(50)
    p1_1.font.bold = True
    p1_1.font.color.rgb = OR_LUMINEUX

    p1_2 = tf1.add_paragraph()
    p1_2.text = "Présentation Institutionnelle"
    p1_2.font.size = Pt(24)
    p1_2.font.bold = True
    p1_2.font.color.rgb = BLANC_PUR

    p1_3 = tf1.add_paragraph()
    p1_3.text = "\n« Former, accompagner et innover pour des territoires durables. »"
    p1_3.font.size = Pt(12)
    p1_3.font.italic = True
    p1_3.font.color.rgb = VERT_DOUX

    # Footer Couverture
    tb_cov_foot = s1.shapes.add_textbox(Inches(0.8), Inches(6.45), Inches(11.7), Inches(0.5))
    pcf = tb_cov_foot.text_frame.paragraphs[0]
    pcf.text = "Présidente Fondatrice : Adama Mbengue  •  Siège : ZAC NORD Thiès Nº2688  •  contact@apatame.com"
    pcf.font.size = Pt(10)
    pcf.font.color.rgb = BLANC_PUR

    # =========================================================================
    # SLIDE 2 : SOMMAIRE DYNAMIQUE (Transition Push Droite)
    # =========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    draw_bg(s2, FOND_PERLE)
    apply_transition(s2, "push", "r")
    add_expert_header(s2, "Sommaire & Table des Matières", "01 / STRUCTURE DU DOCUMENT")

    # Image Droite : Femmes Dynamiques Rayonnantes
    if os.path.exists(img_women):
        s2.shapes.add_picture(img_women, Inches(7.5), Inches(2.1), width=Inches(5.0))

    # Sommaire 8 Piliers
    sommaire_items = [
        ("01", "Présentation & Identité Fondatrice"),
        ("02", "Vision & Ambition Territoriale"),
        ("03", "Mission & Valeurs Fondatrices"),
        ("04", "Nos 6 Domaines d'Intervention"),
        ("05", "Objectifs Stratégiques 2026"),
        ("06", "Notre Approche & Méthode en 5 Étapes"),
        ("07", "Espaces Mutualisés & Chambres d'Hôtes"),
        ("08", "Partenariats, Alliances & Contacts")
    ]

    for idx, (num, label) in enumerate(sommaire_items):
        y = Inches(2.15 + idx * 0.52)
        # Pastille Numéro
        num_circ = s2.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y, Inches(0.42), Inches(0.42))
        num_circ.fill.solid()
        num_circ.fill.fore_color.rgb = BLEU_NUIT if idx % 2 == 0 else VERT_EMERAUDE
        num_circ.line.fill.background()
        np = num_circ.text_frame.paragraphs[0]
        np.text = num
        np.font.size = Pt(10)
        np.font.bold = True
        np.font.color.rgb = BLANC_PUR
        np.alignment = PP_ALIGN.CENTER

        # Libellé
        tb_item = s2.shapes.add_textbox(Inches(1.3), y - Inches(0.05), Inches(5.8), Inches(0.45))
        p_item = tb_item.text_frame.paragraphs[0]
        p_item.text = label
        p_item.font.size = Pt(12)
        p_item.font.bold = True
        p_item.font.color.rgb = NOIR_CHARBON

    add_signature_bottom_bar(s2)

    # =========================================================================
    # SLIDE 3 : QUI SOMMES-NOUS ? (Transition Fade)
    # =========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    draw_bg(s3, FOND_PERLE)
    apply_transition(s3, "fade")
    add_expert_header(s3, "Qui Sommes-Nous ?", "02 / NOTRE IDENTITÉ")

    # Left: Card Définition
    c_qui = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.1), Inches(6.8), Inches(4.3))
    c_qui.fill.solid()
    c_qui.fill.fore_color.rgb = BLANC_PUR
    c_qui.line.color.rgb = BORDER_SUBTILE
    c_qui.line.width = Pt(1)

    qtf = c_qui.text_frame
    qtf.word_wrap = True
    qtf.margin_left = Inches(0.4)
    qtf.margin_top = Inches(0.3)
    qtf.margin_right = Inches(0.4)

    qp1 = qtf.paragraphs[0]
    qp1.text = "Une organisation ancrée au cœur du développement territorial"
    qp1.font.size = Pt(16)
    qp1.font.bold = True
    qp1.font.color.rgb = VERT_FORET

    qp2 = qtf.add_paragraph()
    qp2.text = "\nAPATAM@E est une organisation sénégalaise engagée pour le renforcement des capacités des acteurs locaux, la promotion de l'innovation sociale, la bonne gouvernance et le développement durable.\n\nNotre mission repose sur des approches endogènes et participatives, réconciliant modernité managériale et valeurs positives des terroirs."
    qp2.font.size = Pt(11.5)
    qp2.font.color.rgb = GRIS_TEXTE

    qp3 = qtf.add_paragraph()
    qp3.text = "\nNotre Ambition :"
    qp3.font.size = Pt(13)
    qp3.font.bold = True
    qp3.font.color.rgb = OR_SOLEIL

    qp4 = qtf.add_paragraph()
    qp4.text = "Contribuer à des territoires inclusifs, résilients et prospères."
    qp4.font.size = Pt(11.5)
    qp4.font.bold = True
    qp4.font.color.rgb = BLEU_NUIT

    # Right: Photo Dialogue Aînés & Jeunes
    if os.path.exists(img_dialogue):
        s3.shapes.add_picture(img_dialogue, Inches(7.9), Inches(2.1), width=Inches(4.6))

    add_signature_bottom_bar(s3)

    # =========================================================================
    # SLIDE 4 : NOTRE VISION (Transition Wipe Gauche)
    # =========================================================================
    s4 = prs.slides.add_slide(blank_layout)
    draw_bg(s4, FOND_PERLE)
    apply_transition(s4, "wipe", "l")
    add_expert_header(s4, "Notre Vision Territoriale", "03 / CAP STRATÉGIQUE")

    # Image Droite : Agriculture Solaire Prospère
    if os.path.exists(img_agri):
        s4.shapes.add_picture(img_agri, Inches(6.8), Inches(2.1), width=Inches(5.7))

    # Carte Vision Gauche
    c_vis = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.1), Inches(5.6), Inches(4.3))
    c_vis.fill.solid()
    c_vis.fill.fore_color.rgb = VERT_DOUX
    c_vis.line.color.rgb = VERT_EMERAUDE
    c_vis.line.width = Pt(1.5)

    vtf = c_vis.text_frame
    vtf.word_wrap = True
    vtf.margin_left = Inches(0.4)
    vtf.margin_top = Inches(0.4)
    vtf.margin_right = Inches(0.4)

    vp1 = vtf.paragraphs[0]
    vp1.text = "✦  UNE VISION INCLUSIVE"
    vp1.font.size = Pt(13)
    vp1.font.bold = True
    vp1.font.color.rgb = VERT_EMERAUDE

    vp2 = vtf.add_paragraph()
    vp2.text = "\n« Construire une société sénégalaise égalitaire et inclusive où chaque femme, chaque homme, chaque jeune et chaque personne vulnérable peut réaliser son plein potentiel et contribuer activement au développement durable de son territoire. »"
    vp2.font.size = Pt(15)
    vp2.font.bold = True
    vp2.font.italic = True
    vp2.font.color.rgb = VERT_FORET

    vp3 = vtf.add_paragraph()
    vp3.text = "\n• Justice sociale & égalité des chances\n• Résilience face aux changements climatiques\n• Autonomie économique des communautés locales"
    vp3.font.size = Pt(10.5)
    vp3.font.color.rgb = GRIS_TEXTE

    add_signature_bottom_bar(s4)

    # =========================================================================
    # SLIDE 5 : NOTRE MISSION (Transition Push Haut)
    # =========================================================================
    s5 = prs.slides.add_slide(blank_layout)
    draw_bg(s5, FOND_PERLE)
    apply_transition(s5, "push", "u")
    add_expert_header(s5, "Notre Mission Fondatrice", "04 / ENGAGEMENT D'ACTION")

    # Carte Mission Principale
    c_mis = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.1), Inches(11.73), Inches(1.9))
    c_mis.fill.solid()
    c_mis.fill.fore_color.rgb = BLEU_DOUX
    c_mis.line.color.rgb = BLEU_NUIT
    c_mis.line.width = Pt(1.5)

    mtf = c_mis.text_frame
    mtf.word_wrap = True
    mtf.margin_left = Inches(0.5)
    mtf.margin_top = Inches(0.3)

    mp1 = mtf.paragraphs[0]
    mp1.text = "Promouvoir des actions territoriales inclusives et endogènes garantissant la paix, la sécurité, l'équité et la prospérité durable pour toutes et tous."
    mp1.font.size = Pt(18)
    mp1.font.bold = True
    mp1.font.color.rgb = BLEU_NUIT

    # 4 Cartes d'Impact de la Mission
    mis_piliers = [
        ("🕊️ Paix & Cohésion", "Prévention des conflits et comités locaux de médiation citoyenne.", VERT_EMERAUDE),
        ("⚖️ Équité Sociale", "Accès universel à l'éducation, à la formation et aux ressources.", OR_SOLEIL),
        ("📈 Prospérité Locale", "Modèles économiques circulaires et valorisation des filières agroécologiques.", BLEU_NUIT),
        ("🛡️ Sécurité Citoyenne", "Protection communautaire et renforcement des capacités des acteurs.", VERT_FORET)
    ]
    for idx, (m_t, m_d, m_c) in enumerate(mis_piliers):
        x = Inches(0.8 + idx * 2.95)
        y = Inches(4.25)
        mc = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.8), Inches(2.15))
        mc.fill.solid()
        mc.fill.fore_color.rgb = BLANC_PUR
        mc.line.color.rgb = m_c
        mc.line.width = Pt(1.2)

        mctf = mc.text_frame
        mctf.word_wrap = True
        mctf.margin_top = Inches(0.2)
        mctf.margin_left = Inches(0.2)
        mctf.margin_right = Inches(0.2)

        p1 = mctf.paragraphs[0]
        p1.text = m_t
        p1.font.bold = True
        p1.font.size = Pt(13)
        p1.font.color.rgb = m_c

        p2 = mctf.add_paragraph()
        p2.text = f"\n{m_d}"
        p2.font.size = Pt(10)
        p2.font.color.rgb = GRIS_TEXTE

    add_signature_bottom_bar(s5)

    # =========================================================================
    # SLIDE 6 : NOS VALEURS FONDATRICES (Transition Split Horz)
    # =========================================================================
    s6 = prs.slides.add_slide(blank_layout)
    draw_bg(s6, FOND_PERLE)
    apply_transition(s6, "split")
    add_expert_header(s6, "Nos 6 Valeurs Fondatrices", "05 / SOCLE ÉTHIQUE")

    valeurs = [
        ("01", "ÉQUITÉ", "Nous prônons l'équité et l'égalité des chances.", OR_SOLEIL),
        ("02", "INCLUSION", "Nous plaçons l'humain et les terroirs au cœur de l'action.", BLEU_NUIT),
        ("03", "INNOVATION", "Nous valorisons la créativité et les solutions nouvelles.", OR_SOLEIL),
        ("04", "INTÉGRITÉ", "Nous agissons avec honnêteté, transparence et rigueur.", VERT_EMERAUDE),
        ("05", "COLLABORATION", "Nous croyons à la puissance des synergies multi-acteurs.", BLEU_NUIT),
        ("06", "DURABILITÉ", "Nous bâtissons un avenir résilient et pérenne.", VERT_EMERAUDE)
    ]

    for idx, (num, v_titre, v_desc, v_col) in enumerate(valeurs):
        x = Inches(0.8 + idx * 1.95)
        y = Inches(2.2)
        vc = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(1.85), Inches(4.2))
        vc.fill.solid()
        vc.fill.fore_color.rgb = BLANC_PUR
        vc.line.color.rgb = v_col
        vc.line.width = Pt(1.5)

        vctf = vc.text_frame
        vctf.word_wrap = True
        vctf.margin_top = Inches(0.3)
        vctf.margin_left = Inches(0.15)
        vctf.margin_right = Inches(0.15)

        p1 = vctf.paragraphs[0]
        p1.text = num
        p1.font.bold = True
        p1.font.size = Pt(22)
        p1.font.color.rgb = v_col
        p1.alignment = PP_ALIGN.CENTER

        p2 = vctf.add_paragraph()
        p2.text = f"\n{v_titre}"
        p2.font.bold = True
        p2.font.size = Pt(12)
        p2.font.color.rgb = NOIR_CHARBON
        p2.alignment = PP_ALIGN.CENTER

        p3 = vctf.add_paragraph()
        p3.text = f"\n\n{v_desc}"
        p3.font.size = Pt(10)
        p3.font.color.rgb = GRIS_TEXTE
        p3.alignment = PP_ALIGN.CENTER

    add_signature_bottom_bar(s6)

    # =========================================================================
    # SLIDE 7 : NOS DOMAINES D'INTERVENTION (Transition Fade)
    # =========================================================================
    s7 = prs.slides.add_slide(blank_layout)
    draw_bg(s7, FOND_PERLE)
    apply_transition(s7, "fade")
    add_expert_header(s7, "Nos Domaines d'Intervention", "06 / CHAMPS D'ACTION")

    # Image Droite : Jeunesse et Innovation STEM
    if os.path.exists(img_stem):
        s7.shapes.add_picture(img_stem, Inches(7.5), Inches(2.1), width=Inches(5.0))

    domaines = [
        ("📚  Éducation & Formations Pratiques", "Laboratoires numériques mobiles et bourses d'excellence STEM."),
        ("👩‍💼  Autonomisation Femmes & Jeunes", "Coopératives d'entrepreneuriat solidaire et mentorat."),
        ("♿  Inclusion Sociale & Handicap", "Accessibilité et intégration socio-économique universelle."),
        ("🌱  Résilience Agricole & Eau", "Agroécologie sahélienne et irrigation solaire innovante."),
        ("🏛️  Gouvernance & Participation Citoyenne", "Budget participatif, transparence et contrôle citoyen."),
        ("💡  Innovation Sociale & Économique", "Valorisation des savoirs endogènes et incubation locale.")
    ]

    for idx, (d_titre, d_desc) in enumerate(domaines):
        y = Inches(2.1 + idx * 0.72)
        dc = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(6.4), Inches(0.62))
        dc.fill.solid()
        dc.fill.fore_color.rgb = BLANC_PUR
        dc.line.color.rgb = VERT_EMERAUDE if idx % 2 == 0 else BLEU_NUIT
        dc.line.width = Pt(1.2)

        dctf = dc.text_frame
        dctf.word_wrap = True
        dctf.margin_left = Inches(0.2)
        dctf.margin_top = Inches(0.08)

        p1 = dctf.paragraphs[0]
        p1.text = d_titre
        p1.font.bold = True
        p1.font.size = Pt(11.5)
        p1.font.color.rgb = VERT_FORET

        p2 = dctf.add_paragraph()
        p2.text = d_desc
        p2.font.size = Pt(9.5)
        p2.font.color.rgb = GRIS_TEXTE

    add_signature_bottom_bar(s7)

    # =========================================================================
    # SLIDE 8 : OBJECTIFS STRATÉGIQUES (Transition Zoom)
    # =========================================================================
    s8 = prs.slides.add_slide(blank_layout)
    draw_bg(s8, FOND_PERLE)
    apply_transition(s8, "zoom")
    add_expert_header(s8, "Nos Objectifs Stratégiques 2026", "07 / OBJECTIFS MAJEURS")

    objectifs = [
        ("01", "Éducation & Égalité", "Améliorer l'accès à l'éducation de qualité et à la formation continue pour les filles, les femmes et les jeunes dans les 14 régions.", VERT_EMERAUDE),
        ("02", "Données & Politiques", "Améliorer la production et l'utilisation des résultats de recherche-action pour éclairer et orienter les politiques publiques territoriales.", OR_SOLEIL),
        ("03", "Capacités Institutionnelles", "Renforcer durablement les capacités organisationnelles, financières et opérationnelles des groupements de jeunes et de femmes.", BLEU_NUIT)
    ]

    for idx, (num, obj_t, obj_d, col) in enumerate(objectifs):
        x = Inches(0.8 + idx * 3.95)
        y = Inches(2.2)
        oc = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.7), Inches(4.2))
        oc.fill.solid()
        oc.fill.fore_color.rgb = BLANC_PUR
        oc.line.color.rgb = col
        oc.line.width = Pt(1.5)

        octf = oc.text_frame
        octf.word_wrap = True
        octf.margin_top = Inches(0.4)
        octf.margin_left = Inches(0.3)
        octf.margin_right = Inches(0.3)

        p1 = octf.paragraphs[0]
        p1.text = num
        p1.font.bold = True
        p1.font.size = Pt(36)
        p1.font.color.rgb = col
        p1.alignment = PP_ALIGN.CENTER

        p2 = octf.add_paragraph()
        p2.text = f"\n{obj_t}"
        p2.font.bold = True
        p2.font.size = Pt(15)
        p2.font.color.rgb = NOIR_CHARBON
        p2.alignment = PP_ALIGN.CENTER

        p3 = octf.add_paragraph()
        p3.text = f"\n\n{obj_d}"
        p3.font.size = Pt(11)
        p3.font.color.rgb = GRIS_TEXTE
        p3.alignment = PP_ALIGN.CENTER

    add_signature_bottom_bar(s8)

    # =========================================================================
    # SLIDE 9 : NOTRE APPROCHE SYSTÉMIQUE (Transition Wheel 4 Rayons)
    # =========================================================================
    s9 = prs.slides.add_slide(blank_layout)
    draw_bg(s9, FOND_PERLE)
    apply_transition(s9, "wheel")
    add_expert_header(s9, "Notre Approche Systémique & Circulaire", "08 / MODÈLE D'ACTION")

    approches = [
        ("01. FORMER", "Renforcer les compétences techniques, managériales et civiques des acteurs territoriaux.", VERT_EMERAUDE),
        ("02. ACCOMPAGNER", "Fournir un accompagnement sur-mesure sur le terrain à chaque étape du déploiement.", OR_SOLEIL),
        ("03. INNOVER", "Co-concevoir des solutions durables adaptées aux spécificités culturelles et écologiques.", BLEU_NUIT),
        ("04. TRANSFORMER", "Créer un impact systémique mesurable et pérenne sur la vie des populations locales.", VERT_FORET)
    ]

    for idx, (ap_t, ap_d, ap_c) in enumerate(approches):
        x = Inches(0.8 + idx * 2.95)
        y = Inches(2.2)
        ac = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.8), Inches(4.2))
        ac.fill.solid()
        ac.fill.fore_color.rgb = BLANC_PUR
        ac.line.color.rgb = ap_c
        ac.line.width = Pt(1.5)

        actf = ac.text_frame
        actf.word_wrap = True
        actf.margin_top = Inches(0.4)
        actf.margin_left = Inches(0.2)
        actf.margin_right = Inches(0.2)

        p1 = actf.paragraphs[0]
        p1.text = ap_t
        p1.font.bold = True
        p1.font.size = Pt(14)
        p1.font.color.rgb = ap_c
        p1.alignment = PP_ALIGN.CENTER

        p2 = actf.add_paragraph()
        p2.text = f"\n\n{ap_d}"
        p2.font.size = Pt(11)
        p2.font.color.rgb = GRIS_TEXTE
        p2.alignment = PP_ALIGN.CENTER

    add_signature_bottom_bar(s9)

    # =========================================================================
    # SLIDE 10 : NOS SERVICES (Transition Push Droite)
    # =========================================================================
    s10 = prs.slides.add_slide(blank_layout)
    draw_bg(s10, FOND_PERLE)
    apply_transition(s10, "push", "r")
    add_expert_header(s10, "Notre Offre de Services", "09 / EXPERTISE OPÉRATIONNELLE")

    services = [
        ("Formation Professionnelle", "Ingénierie pédagogique, sessions certifiantes et ateliers pratiques.", BLEU_NUIT),
        ("Conseil Stratégique", "Plans de développement communal, gouvernance locale et conduite du changement.", VERT_EMERAUDE),
        ("Recherche & Observatoire", "Études d'impact, diagnostics territoriaux et collecte de données géolocalisées.", OR_SOLEIL),
        ("Accompagnement Terrain", "Coaching individuel et collectif des coopératives et porteurs de projet.", VERT_EMERAUDE),
        ("Études & Diagnostics", "Audits institutionnels, faisabilité technique et évaluations environnementales.", OR_SOLEIL),
        ("Renforcement des Capacités", "Leadership communautaire, gestion financière et plaidoyer citoyen.", BLEU_NUIT)
    ]

    for idx, (s_title, s_desc, s_col) in enumerate(services):
        row = idx // 3
        col = idx % 3
        x = Inches(0.8 + col * 4.0)
        y = Inches(2.2 + row * 2.1)
        sc = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.75), Inches(1.85))
        sc.fill.solid()
        sc.fill.fore_color.rgb = BLANC_PUR
        sc.line.color.rgb = s_col
        sc.line.width = Pt(1.5)

        sctf = sc.text_frame
        sctf.word_wrap = True
        sctf.margin_top = Inches(0.2)
        sctf.margin_left = Inches(0.25)

        p1 = sctf.paragraphs[0]
        p1.text = f"✦  {s_title}"
        p1.font.bold = True
        p1.font.size = Pt(13)
        p1.font.color.rgb = s_col

        p2 = sctf.add_paragraph()
        p2.text = f"\n{s_desc}"
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = GRIS_TEXTE

    add_signature_bottom_bar(s10)

    # =========================================================================
    # SLIDE 11 : NOTRE MÉTHODE D'INTERVENTION (Transition Wipe Bas)
    # =========================================================================
    s11 = prs.slides.add_slide(blank_layout)
    draw_bg(s11, FOND_PERLE)
    apply_transition(s11, "wipe", "d")
    add_expert_header(s11, "Notre Méthode d'Intervention en 5 Étapes", "10 / DÉMARCHE QUALITÉ")

    method_steps = [
        ("01", "DIAGNOSTIC", "Analyse fine des besoins et écoute active des bénéficiaires sur le terrain."),
        ("02", "CO-CONSTRUCTION", "Définition conjointe des plans d'action avec les communautés locales."),
        ("03", "MISE EN ŒUVRE", "Exécution agile, déploiement des formations et des infrastructures."),
        ("04", "SUIVI & ÉVALUATION", "Mesure rigoureuse des indicateurs de performance et d'impact réel."),
        ("05", "CAPITALISATION", "Modélisation des bonnes pratiques et dissémination à grande échelle.")
    ]

    for idx, (num, step_title, step_desc) in enumerate(method_steps):
        x = Inches(0.8 + idx * 2.35)
        y = Inches(2.2)
        mc = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.2), Inches(4.2))
        mc.fill.solid()
        mc.fill.fore_color.rgb = BLANC_PUR
        mc.line.color.rgb = OR_SOLEIL if idx == 2 else (VERT_EMERAUDE if idx % 2 == 0 else BLEU_NUIT)
        mc.line.width = Pt(1.5)

        mctf = mc.text_frame
        mctf.word_wrap = True
        mctf.margin_top = Inches(0.3)
        mctf.margin_left = Inches(0.15)
        mctf.margin_right = Inches(0.15)

        p1 = mctf.paragraphs[0]
        p1.text = num
        p1.font.bold = True
        p1.font.size = Pt(24)
        p1.font.color.rgb = OR_SOLEIL
        p1.alignment = PP_ALIGN.CENTER

        p2 = mctf.add_paragraph()
        p2.text = f"\n{step_title}"
        p2.font.bold = True
        p2.font.size = Pt(11.5)
        p2.font.color.rgb = NOIR_CHARBON
        p2.alignment = PP_ALIGN.CENTER

        p3 = mctf.add_paragraph()
        p3.text = f"\n\n{step_desc}"
        p3.font.size = Pt(10)
        p3.font.color.rgb = GRIS_TEXTE
        p3.alignment = PP_ALIGN.CENTER

    add_signature_bottom_bar(s11)

    # =========================================================================
    # SLIDE 12 : ESPACES MUTUALISÉS & CHAMBRES D'HÔTES (Transition Fade Immersif)
    # =========================================================================
    s12 = prs.slides.add_slide(blank_layout)
    draw_bg(s12, FOND_PERLE)
    apply_transition(s12, "fade")
    add_expert_header(s12, "Espaces Mutualisés & Chambres d'Hôtes", "11 / INFRASTRUCTURES D'ACCUEIL")

    # Gauche : Photo Salles de Réunion & Séminaire
    if os.path.exists(img_seminaire):
        s12.shapes.add_picture(img_seminaire, Inches(0.8), Inches(2.15), width=Inches(5.6))
    
    c_salle = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.7), Inches(5.6), Inches(1.75))
    c_salle.fill.solid()
    c_salle.fill.fore_color.rgb = BLANC_PUR
    c_salle.line.color.rgb = VERT_EMERAUDE
    c_salle.line.width = Pt(1.5)
    stf = c_salle.text_frame
    stf.word_wrap = True
    stf.margin_left = Inches(0.25)
    stf.margin_top = Inches(0.15)
    sp1 = stf.paragraphs[0]
    sp1.text = "🏛️ Salles de Séminaires (ZAC Nord Thiès)"
    sp1.font.bold = True
    sp1.font.size = Pt(12)
    sp1.font.color.rgb = VERT_FORET
    sp2 = stf.add_paragraph()
    sp2.text = "• Climatisation, Vidéoprojecteur interactif, Wifi très haut débit.\n• Idéal pour Séminaires, Assemblées Générales et Ateliers de formation."
    sp2.font.size = Pt(9.5)
    sp2.font.color.rgb = GRIS_TEXTE

    # Droite : Photo Chambres d'Hôtes Saly & Thiès
    if os.path.exists(img_chambre):
        s12.shapes.add_picture(img_chambre, Inches(6.8), Inches(2.15), width=Inches(5.7))

    c_chambre = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(4.7), Inches(5.7), Inches(1.75))
    c_chambre.fill.solid()
    c_chambre.fill.fore_color.rgb = BLANC_PUR
    c_chambre.line.color.rgb = OR_SOLEIL
    c_chambre.line.width = Pt(1.5)
    htf = c_chambre.text_frame
    htf.word_wrap = True
    htf.margin_left = Inches(0.25)
    htf.margin_top = Inches(0.15)
    hp1 = htf.paragraphs[0]
    hp1.text = "🏡 Chambres d'Hôtes (Thiès & Saly Portudal)"
    hp1.font.bold = True
    hp1.font.size = Pt(12)
    hp1.font.color.rgb = OR_SOLEIL
    hp2 = htf.add_paragraph()
    hp2.text = "• Résidence Thiès (missions pro) & Villa Saly (piscine & retraites de travail).\n• Petit-déjeuner bio local inclus, conciergerie et transfert aéroport AIBD."
    hp2.font.size = Pt(9.5)
    hp2.font.color.rgb = GRIS_TEXTE

    add_signature_bottom_bar(s12)

    # =========================================================================
    # SLIDE 13 : ALLIANCES & APPEL À PARTENARIAT (Transition Push Droite)
    # =========================================================================
    s13 = prs.slides.add_slide(blank_layout)
    draw_bg(s13, FOND_PERLE)
    apply_transition(s13, "push", "r")
    add_expert_header(s13, "Construisons Ensemble des Territoires Durables", "12 / APPEL À PARTENARIAT")

    # Grand Cadre Central
    p_box = s13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.15), Inches(11.73), Inches(1.6))
    p_box.fill.solid()
    p_box.fill.fore_color.rgb = VERT_DOUX
    p_box.line.color.rgb = VERT_EMERAUDE
    p_box.line.width = Pt(1.5)

    ptf = p_box.text_frame
    ptf.word_wrap = True
    ptf.margin_left = Inches(0.4)
    ptf.margin_top = Inches(0.25)

    pp1 = ptf.paragraphs[0]
    pp1.text = "Un écosystème ouvert et multi-acteurs pour maximiser l'impact"
    pp1.font.bold = True
    pp1.font.size = Pt(16)
    pp1.font.color.rgb = VERT_FORET

    pp2 = ptf.add_paragraph()
    pp2.text = "APATAM@E co-construit des solutions sur-mesure avec les institutions de l'État, les collectivités territoriales, les partenaires techniques et financiers, le secteur privé et la société civile."
    pp2.font.size = Pt(11)
    pp2.font.color.rgb = NOIR_CHARBON

    # 5 Cartes de Partenariat
    partenaires_cats = [
        ("🏛️ Institutions", "Ministères & Agences de l'État"),
        ("🌍 Collectivités", "Mairies & Conseils Départementaux"),
        ("🤝 Bailleurs & PTF", "Coopération & Fonds Internationaux"),
        ("💼 Secteur Privé", "Entreprises RSE & Fondations"),
        ("👥 Société Civile", "Associations & Mouvements Jeunesse")
    ]
    for idx, (cat_t, cat_sub) in enumerate(partenaires_cats):
        x = Inches(0.8 + idx * 2.35)
        y = Inches(4.0)
        c_part = s13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.2), Inches(2.35))
        c_part.fill.solid()
        c_part.fill.fore_color.rgb = BLANC_PUR
        c_part.line.color.rgb = BORDER_SUBTILE

        cptf = c_part.text_frame
        cptf.word_wrap = True
        cptf.margin_top = Inches(0.4)
        cptf.margin_left = Inches(0.15)
        cptf.margin_right = Inches(0.15)

        p1 = cptf.paragraphs[0]
        p1.text = cat_t
        p1.font.bold = True
        p1.font.size = Pt(12)
        p1.font.color.rgb = BLEU_NUIT
        p1.alignment = PP_ALIGN.CENTER

        p2 = cptf.add_paragraph()
        p2.text = f"\n{cat_sub}"
        p2.font.size = Pt(9.5)
        p2.font.color.rgb = GRIS_TEXTE
        p2.alignment = PP_ALIGN.CENTER

    add_signature_bottom_bar(s13)

    # =========================================================================
    # SLIDE 14 : CONTACTS & COORDONNÉES OFFICIELLES (Transition Fade)
    # =========================================================================
    s14 = prs.slides.add_slide(blank_layout)
    draw_bg(s14, FOND_PERLE)
    apply_transition(s14, "fade")
    add_expert_header(s14, "Coordonnées Officielles & Contact", "13 / NOUS CONTACTER")

    # Carte Contact
    c_contact = s14.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.15), Inches(7.5), Inches(4.3))
    c_contact.fill.solid()
    c_contact.fill.fore_color.rgb = BLANC_PUR
    c_contact.line.color.rgb = VERT_EMERAUDE
    c_contact.line.width = Pt(1.5)

    ctf = c_contact.text_frame
    ctf.word_wrap = True
    ctf.margin_left = Inches(0.4)
    ctf.margin_top = Inches(0.3)

    cp1 = ctf.paragraphs[0]
    cp1.text = "APATAM@E — Direction Générale"
    cp1.font.bold = True
    cp1.font.size = Pt(16)
    cp1.font.color.rgb = BLEU_NUIT

    cp2 = ctf.add_paragraph()
    cp2.text = "Présidente Fondatrice : Adama Mbengue\n"
    cp2.font.size = Pt(12)
    cp2.font.bold = True
    cp2.font.color.rgb = VERT_FORET

    cp3 = ctf.add_paragraph()
    cp3.text = "📍 Siège social : ZAC NORD Thiès Nº2688, Rond Point Mosquée, Thiès - Sénégal\n"
    cp3.font.size = Pt(11)
    cp3.font.color.rgb = NOIR_CHARBON

    cp4 = ctf.add_paragraph()
    cp4.text = "📞 Bureau : +221 33 999 28 52   |   Portable : +221 77 510 20 38\n"
    cp4.font.size = Pt(11)
    cp4.font.color.rgb = NOIR_CHARBON

    cp5 = ctf.add_paragraph()
    cp5.text = "📧 E-mails : contact@apatame.com   |   secretariat@apatame.org\n"
    cp5.font.size = Pt(11)
    cp5.font.color.rgb = NOIR_CHARBON

    cp6 = ctf.add_paragraph()
    cp6.text = "🌐 Site Web : www.apatame.org (ou http://localhost:5173)\n🏡 Espaces Mutualisés & Chambres d'Hôtes : Thiès & Saly"
    cp6.font.size = Pt(11)
    cp6.font.bold = True
    cp6.font.color.rgb = OR_SOLEIL

    # Photo Présidente à Droite
    if os.path.exists(img_pres):
        s14.shapes.add_picture(img_pres, Inches(8.7), Inches(2.15), width=Inches(3.8))

    add_signature_bottom_bar(s14)

    # =========================================================================
    # SLIDE 15 : MERCI & CLÔTURE INSPIRANTE (Transition Zoom Puissante)
    # =========================================================================
    s15 = prs.slides.add_slide(blank_layout)
    draw_bg(s15, VERT_FORET)
    apply_transition(s15, "zoom")

    # Image Droite : Pousse de plante lumineuse entre les mains
    if os.path.exists(img_sprout):
        s15.shapes.add_picture(img_sprout, Inches(6.6), Inches(0.8), width=Inches(6.0))

    # Logo Gauche
    if os.path.exists(logo_path):
        s15.shapes.add_picture(logo_path, Inches(0.9), Inches(1.0), height=Inches(1.5))

    tb_merci = s15.shapes.add_textbox(Inches(0.9), Inches(2.8), Inches(5.5), Inches(3.5))
    mtf = tb_merci.text_frame
    mtf.word_wrap = True

    mp1 = mtf.paragraphs[0]
    mp1.text = "MERCI"
    mp1.font.size = Pt(54)
    mp1.font.bold = True
    mp1.font.color.rgb = OR_LUMINEUX

    mp2 = mtf.add_paragraph()
    mp2.text = "POUR VOTRE ATTENTION."
    mp2.font.size = Pt(22)
    mp2.font.bold = True
    mp2.font.color.rgb = BLANC_PUR

    mp3 = mtf.add_paragraph()
    mp3.text = "\nEnsemble, bâtissons des territoires durables et prospères pour les générations futures."
    mp3.font.size = Pt(13)
    mp3.font.italic = True
    mp3.font.color.rgb = VERT_DOUX

    mp4 = mtf.add_paragraph()
    mp4.text = "\nwww.apatame.org  •  contact@apatame.com"
    mp4.font.size = Pt(11)
    mp4.font.bold = True
    mp4.font.color.rgb = BLANC_PUR

    # Output filenames
    filename = "APATAM_E_Presentation_Expert_Animee_2026.pptx"
    
    # 1. Parent folder
    parent_dir = os.path.abspath(os.path.join(base_dir, ".."))
    parent_path = os.path.join(parent_dir, filename)
    prs.save(parent_path)
    print(f"[OK] Enregistré dans le dossier parent : {parent_path}")

    # 2. Downloads folder
    user_home = os.path.expanduser("~")
    downloads_path = os.path.join(user_home, "Downloads", filename)
    prs.save(downloads_path)
    print(f"[OK] Enregistré dans les Téléchargements : {downloads_path}")

    # 3. Project folder
    project_path = os.path.join(base_dir, filename)
    prs.save(project_path)
    print(f"[OK] Enregistré dans le projet : {project_path}")

    # 4. E:\noppal\apatamae folder
    if os.path.exists(e_drive_dir):
        e_path = os.path.join(e_drive_dir, filename)
        prs.save(e_path)
        print(f"[OK] Enregistré dans E:\\noppal\\apatamae : {e_path}")

if __name__ == "__main__":
    create_expert_presentation()
