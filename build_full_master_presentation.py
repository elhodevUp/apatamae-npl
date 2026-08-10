import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml import parse_xml

def create_master_presentation():
    prs = Presentation()
    # Format 16:9 Widescreen HD
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # --- PALETTE DE COULEURS OFFICIELLES APATAM@E ---
    BLEU_INSTITUTIONNEL = RGBColor(10, 34, 64)    # #0A2240
    VERT_TERRITOIRE     = RGBColor(18, 102, 53)   # #126635
    VERT_FORET          = RGBColor(1, 45, 29)     # #012D1D
    OR_RAYONNEMENT      = RGBColor(222, 154, 25)  # #DE9A19
    OR_CLAIR            = RGBColor(245, 190, 44)  # #F5BE2C
    GRIS_FOND           = RGBColor(250, 251, 252) # #FAFBFC
    BLANC               = RGBColor(255, 255, 255)
    NOIR_TEXTE          = RGBColor(25, 28, 29)
    GRIS_TEXTE          = RGBColor(80, 90, 85)
    BORDER_GRIS         = RGBColor(225, 230, 228)
    VERT_PALE           = RGBColor(230, 242, 235)
    BLEU_PALE           = RGBColor(232, 240, 250)
    OR_PALE             = RGBColor(254, 247, 230)

    base_dir = os.path.dirname(os.path.abspath(__file__))
    artifacts_dir = r"C:\Users\LENOVO E15\.gemini\antigravity-ide\brain\2584ad8b-cb21-4a7f-a2cd-17da14691d98"
    e_drive_dir = r"E:\noppal\apatamae"

    # Logo
    logo_path = os.path.join(base_dir, "public", "logo.png")
    if not os.path.exists(logo_path):
        logo_path = os.path.join(e_drive_dir, "logo.png.jpeg")

    # Images from E: drive & artifacts
    img_cover = os.path.join(e_drive_dir, "WhatsApp Image 2026-08-04 at 16.37.15.jpeg")
    img_sommaire = os.path.join(e_drive_dir, "WhatsApp Image 2026-08-04 at 16.37.15 (1).jpeg")
    img_baobab = os.path.join(e_drive_dir, "WhatsApp Image 2026-08-04 at 16.37.15 (2).jpeg")
    img_mission = os.path.join(e_drive_dir, "WhatsApp Image 2026-08-04 at 16.37.16.jpeg")
    img_valeurs = os.path.join(e_drive_dir, "WhatsApp Image 2026-08-04 at 16.37.16 (1).jpeg")
    img_domaines = os.path.join(e_drive_dir, "WhatsApp Image 2026-08-04 at 16.37.16 (2).jpeg")
    img_objectifs = os.path.join(e_drive_dir, "WhatsApp Image 2026-08-04 at 16.37.17.jpeg")
    img_merci = os.path.join(e_drive_dir, "WhatsApp Image 2026-08-04 at 16.37.17 (1).jpeg")
    img_vision = os.path.join(e_drive_dir, "WhatsApp Image 2026-08-04 at 16.37.17 (2).jpeg")
    img_pres = os.path.join(base_dir, "public", "images", "professional_corporate_portrait_of_a_west_african_woman_leader_elegant_and_screen.png")
    img_seminaire = os.path.join(artifacts_dir, "pptx_salles_seminaire_1785866596542.png")
    img_chambre = os.path.join(artifacts_dir, "pptx_chambre_hotes_1785866613303.png")

    def apply_transition(slide, trans_type="fade"):
        """Injecte les animations de transition natives PowerPoint"""
        try:
            xml = f'<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:{trans_type}/></p:transition>'
            slide._element.append(parse_xml(xml))
        except Exception as e:
            pass

    def draw_bg(slide, color=GRIS_FOND):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = color
        bg.line.fill.background()
        return bg

    def add_official_header(slide, title_text, category_text="APATAM@E • INSTITUTIONAL", dark_mode=False):
        """Header officiel avec logo, titre grande écriture, slogan et ligne dorée décorative"""
        # Logo en haut à gauche
        if os.path.exists(logo_path):
            slide.shapes.add_picture(logo_path, Inches(0.8), Inches(0.4), height=Inches(0.9))

        # Titre & Slogan
        tb = slide.shapes.add_textbox(Inches(2.2), Inches(0.35), Inches(9.5), Inches(0.9))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p0 = tf.paragraphs[0]
        p0.text = "APATAM@E"
        p0.font.size = Pt(16)
        p0.font.bold = True
        p0.font.color.rgb = BLEU_INSTITUTIONNEL if not dark_mode else OR_CLAIR

        p_slogan = tf.add_paragraph()
        p_slogan.text = "Former, accompagner et innover pour des territoires durables."
        p_slogan.font.size = Pt(10)
        p_slogan.font.italic = True
        p_slogan.font.color.rgb = VERT_TERRITOIRE if not dark_mode else BLANC

        # Titre Principal de la Slide (Grandes Écritures)
        tb_titre = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.7))
        p_t = tb_titre.text_frame.paragraphs[0]
        p_t.text = title_text
        p_t.font.size = Pt(26)
        p_t.font.bold = True
        p_t.font.color.rgb = VERT_FORET if not dark_mode else BLANC

        # Ligne Dorée Décorative sous le titre
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.95), Inches(11.7), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = OR_RAYONNEMENT
        line.line.fill.background()

    def add_bottom_bar(slide, items=[("Inclusion", VERT_TERRITOIRE), ("Innovation", OR_RAYONNEMENT), ("Territoires", BLEU_INSTITUTIONNEL), ("Durabilité", VERT_TERRITOIRE)]):
        """Barre inférieure stylisée avec arrondis et valeurs clés"""
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.55))
        bar.fill.solid()
        bar.fill.fore_color.rgb = BLEU_INSTITUTIONNEL
        bar.line.fill.background()
        
        btf = bar.text_frame
        btf.word_wrap = True
        p = btf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        
        text_str = "   •   ".join([label for label, _ in items])
        p.text = f"✦  {text_str}  ✦"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = BLANC

    # =========================================================================
    # SLIDE 1 : COUVERTURE OFFICIELLE (PRÉSENTATION INSTITUTIONNELLE)
    # =========================================================================
    s1 = prs.slides.add_slide(blank_layout)
    draw_bg(s1, BLANC)
    apply_transition(s1, "fade")

    # Si l'image de couverture existe, on l'utilise directement en pleine page
    if os.path.exists(img_cover):
        s1.shapes.add_picture(img_cover, 0, 0, width=Inches(13.333), height=Inches(7.5))
    else:
        # Fallback design élégant
        draw_bg(s1, BLANC)
        if os.path.exists(logo_path):
            s1.shapes.add_picture(logo_path, Inches(1.0), Inches(1.0), height=Inches(1.5))
        
        tb = s1.shapes.add_textbox(Inches(1.0), Inches(3.0), Inches(8.0), Inches(2.5))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = "PRÉSENTATION INSTITUTIONNELLE"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = BLEU_INSTITUTIONNEL

        p2 = tf.add_paragraph()
        p2.text = "Former, accompagner et innover pour des territoires durables."
        p2.font.size = Pt(16)
        p2.font.italic = True
        p2.font.color.rgb = VERT_TERRITOIRE

    # =========================================================================
    # SLIDE 2 : SOMMAIRE (10 POINTS)
    # =========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    draw_bg(s2, BLANC)
    apply_transition(s2, "push")

    if os.path.exists(img_sommaire):
        s2.shapes.add_picture(img_sommaire, 0, 0, width=Inches(13.333), height=Inches(7.5))
    else:
        add_official_header(s2, "SOMMAIRE")
        sommaire_list = [
            ("01", "Présentation d'APATAM@E"),
            ("02", "Vision"),
            ("03", "Mission"),
            ("04", "Nos Valeurs"),
            ("05", "Domaines d'intervention"),
            ("06", "Objectifs stratégiques"),
            ("07", "Projets et initiatives"),
            ("08", "Résultats attendus"),
            ("09", "Partenariats & Alliances"),
            ("10", "Contacts & Réservations")
        ]
        for idx, (num, label) in enumerate(sommaire_list):
            row = idx % 5
            col = idx // 5
            x = Inches(0.8 + col * 4.5)
            y = Inches(2.2 + row * 0.8)

            c_num = s2.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.55), Inches(0.55))
            c_num.fill.solid()
            c_num.fill.fore_color.rgb = BLEU_INSTITUTIONNEL if col == 0 else VERT_TERRITOIRE
            c_num.line.fill.background()
            p = c_num.text_frame.paragraphs[0]
            p.text = num
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = BLANC
            p.alignment = PP_ALIGN.CENTER

            tb_lbl = s2.shapes.add_textbox(x + Inches(0.65), y, Inches(3.6), Inches(0.55))
            p_l = tb_lbl.text_frame.paragraphs[0]
            p_l.text = label
            p_l.font.size = Pt(13)
            p_l.font.bold = True
            p_l.font.color.rgb = NOIR_TEXTE

    # =========================================================================
    # SLIDE 3 : QUI SOMMES-NOUS ?
    # =========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    draw_bg(s3, BLANC)
    apply_transition(s3, "fade")

    if os.path.exists(img_baobab):
        s3.shapes.add_picture(img_baobab, 0, 0, width=Inches(13.333), height=Inches(7.5))
    else:
        add_official_header(s3, "QUI SOMMES-NOUS ?")
        # Carte descriptif
        c_desc = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.2), Inches(6.5), Inches(4.2))
        c_desc.fill.solid()
        c_desc.fill.fore_color.rgb = GRIS_FOND
        c_desc.line.color.rgb = BORDER_GRIS
        dtf = c_desc.text_frame
        dtf.word_wrap = True
        dtf.margin_left = Inches(0.4)
        dtf.margin_top = Inches(0.3)
        
        p1 = dtf.paragraphs[0]
        p1.text = "Une organisation dédiée à l'action territoriale"
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = BLEU_INSTITUTIONNEL

        p2 = dtf.add_paragraph()
        p2.text = "\nAPATAM@E est une organisation sénégalaise qui œuvre pour le renforcement des capacités des acteurs territoriaux, la promotion de l'innovation sociale, la gouvernance locale et le développement durable à travers des approches inclusives et endogènes.\n\nNotre ambition :\nContribuer à des territoires inclusifs, résilients et prospères."
        p2.font.size = Pt(12)
        p2.font.color.rgb = NOIR_TEXTE

        add_bottom_bar(s3)

    # =========================================================================
    # SLIDE 4 : VISION
    # =========================================================================
    s4 = prs.slides.add_slide(blank_layout)
    draw_bg(s4, BLANC)
    apply_transition(s4, "wipe")

    if os.path.exists(img_vision):
        s4.shapes.add_picture(img_vision, 0, 0, width=Inches(13.333), height=Inches(7.5))
    else:
        add_official_header(s4, "NOTRE VISION")
        # Vision Box
        vbox = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.2), Inches(7.5), Inches(4.0))
        vbox.fill.solid()
        vbox.fill.fore_color.rgb = VERT_PALE
        vbox.line.color.rgb = VERT_TERRITOIRE
        vbox.line.width = Pt(1.5)
        vtf = vbox.text_frame
        vtf.word_wrap = True
        vtf.margin_left = Inches(0.5)
        vtf.margin_top = Inches(0.4)

        vp1 = vtf.paragraphs[0]
        vp1.text = "Construire une société sénégalaise égalitaire et inclusive où chaque femme, chaque homme, chaque jeune et chaque personne vulnérable peut réaliser son potentiel et contribuer au développement durable des territoires."
        vp1.font.size = Pt(18)
        vp1.font.bold = True
        vp1.font.color.rgb = VERT_FORET

        add_bottom_bar(s4, [("Paix", OR_RAYONNEMENT), ("Équité", BLEU_INSTITUTIONNEL), ("Prospérité", VERT_TERRITOIRE), ("Sécurité", VERT_TERRITOIRE)])

    # =========================================================================
    # SLIDE 5 : MISSION
    # =========================================================================
    s5 = prs.slides.add_slide(blank_layout)
    draw_bg(s5, BLANC)
    apply_transition(s5, "push")

    if os.path.exists(img_mission):
        s5.shapes.add_picture(img_mission, 0, 0, width=Inches(13.333), height=Inches(7.5))
    else:
        add_official_header(s5, "NOTRE MISSION")
        # Mission Box
        mbox = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.2), Inches(7.5), Inches(4.0))
        mbox.fill.solid()
        mbox.fill.fore_color.rgb = BLEU_PALE
        mbox.line.color.rgb = BLEU_INSTITUTIONNEL
        mbox.line.width = Pt(1.5)
        mtf = mbox.text_frame
        mtf.word_wrap = True
        mtf.margin_left = Inches(0.5)
        mtf.margin_top = Inches(0.4)

        mp1 = mtf.paragraphs[0]
        mp1.text = "Promouvoir des actions territoriales inclusives et endogènes garantissant la paix, la sécurité, l'équité et la prospérité durable pour toutes et tous."
        mp1.font.size = Pt(18)
        mp1.font.bold = True
        mp1.font.color.rgb = BLEU_INSTITUTIONNEL

        add_bottom_bar(s5, [("Paix", OR_RAYONNEMENT), ("Équité", BLEU_INSTITUTIONNEL), ("Prospérité", VERT_TERRITOIRE), ("Sécurité", VERT_TERRITOIRE)])

    # =========================================================================
    # SLIDE 6 : NOS VALEURS FONDAMENTALES (6 VALEURS)
    # =========================================================================
    s6 = prs.slides.add_slide(blank_layout)
    draw_bg(s6, BLANC)
    apply_transition(s6, "fade")

    if os.path.exists(img_valeurs):
        s6.shapes.add_picture(img_valeurs, 0, 0, width=Inches(13.333), height=Inches(7.5))
    else:
        add_official_header(s6, "NOS VALEURS FONDAMENTALES")
        valeurs_cards = [
            ("ÉQUITÉ", "Nous prônons l'équité et l'égalité des chances.", OR_RAYONNEMENT),
            ("INCLUSION", "Nous plaçons l'humain au cœur de nos actions.", BLEU_INSTITUTIONNEL),
            ("INNOVATION", "Nous valorisons la créativité et les solutions nouvelles.", OR_RAYONNEMENT),
            ("INTÉGRITÉ", "Nous agissons avec honnêteté, transparence et responsabilité.", VERT_TERRITOIRE),
            ("COLLABORATION", "Nous croyons à la force du partenariat.", BLEU_INSTITUTIONNEL),
            ("DURABILITÉ", "Nous nous engageons pour un avenir durable et résilient.", VERT_TERRITOIRE)
        ]
        for i, (v_titre, v_desc, v_col) in enumerate(valeurs_cards):
            x = Inches(0.8 + i * 1.95)
            y = Inches(2.2)
            vc = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(1.8), Inches(4.0))
            vc.fill.solid()
            vc.fill.fore_color.rgb = GRIS_FOND
            vc.line.color.rgb = v_col
            vc.line.width = Pt(1.5)
            
            vctf = vc.text_frame
            vctf.word_wrap = True
            vctf.margin_top = Inches(0.4)
            vctf.margin_left = Inches(0.15)
            vctf.margin_right = Inches(0.15)

            p1 = vctf.paragraphs[0]
            p1.text = v_titre
            p1.font.bold = True
            p1.font.size = Pt(13)
            p1.font.color.rgb = v_col
            p1.alignment = PP_ALIGN.CENTER

            p2 = vctf.add_paragraph()
            p2.text = f"\n\n{v_desc}"
            p2.font.size = Pt(10)
            p2.font.color.rgb = NOIR_TEXTE
            p2.alignment = PP_ALIGN.CENTER

    # =========================================================================
    # SLIDE 7 : DOMAINES D'INTERVENTION
    # =========================================================================
    s7 = prs.slides.add_slide(blank_layout)
    draw_bg(s7, BLANC)
    apply_transition(s7, "wipe")

    if os.path.exists(img_domaines):
        s7.shapes.add_picture(img_domaines, 0, 0, width=Inches(13.333), height=Inches(7.5))
    else:
        add_official_header(s7, "DOMAINES D'INTERVENTION")
        domaines_list = [
            "Éducation et formation",
            "Autonomisation des femmes et des jeunes",
            "Inclusion sociale et handicap",
            "Résilience agricole et sécurité alimentaire",
            "Gouvernance locale et participation citoyenne",
            "Innovation sociale et entrepreneuriat"
        ]
        for i, dom in enumerate(domaines_list):
            y = Inches(2.2 + i * 0.7)
            c = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(6.5), Inches(0.6))
            c.fill.solid()
            c.fill.fore_color.rgb = VERT_PALE
            c.line.color.rgb = VERT_TERRITOIRE
            p = c.text_frame.paragraphs[0]
            p.text = f"✦  {dom}"
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = VERT_FORET

    # =========================================================================
    # SLIDE 8 : OBJECTIFS STRATÉGIQUES
    # =========================================================================
    s8 = prs.slides.add_slide(blank_layout)
    draw_bg(s8, BLANC)
    apply_transition(s8, "push")

    if os.path.exists(img_objectifs):
        s8.shapes.add_picture(img_objectifs, 0, 0, width=Inches(13.333), height=Inches(7.5))
    else:
        add_official_header(s8, "OBJECTIFS STRATÉGIQUES")
        objs = [
            ("01", "Améliorer l'accès à l'éducation et à la formation pour les filles, les femmes et les jeunes.", VERT_TERRITOIRE),
            ("02", "Améliorer la production et l'utilisation des résultats pour éclairer les politiques publiques.", OR_RAYONNEMENT),
            ("03", "Renforcer les capacités institutionnelles des organisations de jeunes et de femmes.", BLEU_INSTITUTIONNEL)
        ]
        for i, (num, desc, col) in enumerate(objs):
            x = Inches(0.8 + i * 3.95)
            y = Inches(2.2)
            oc = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.7), Inches(4.0))
            oc.fill.solid()
            oc.fill.fore_color.rgb = GRIS_FOND
            oc.line.color.rgb = col
            oc.line.width = Pt(1.5)
            
            octf = oc.text_frame
            octf.word_wrap = True
            octf.margin_top = Inches(0.4)
            octf.margin_left = Inches(0.3)
            octf.margin_right = Inches(0.3)

            p1 = octf.paragraphs[0]
            p1.text = num
            p1.font.bold = True
            p1.font.size = Pt(28)
            p1.font.color.rgb = col
            p1.alignment = PP_ALIGN.CENTER

            p2 = octf.add_paragraph()
            p2.text = f"\n\n{desc}"
            p2.font.size = Pt(12)
            p2.font.color.rgb = NOIR_TEXTE
            p2.alignment = PP_ALIGN.CENTER

    # =========================================================================
    # SLIDE 9 : NOTRE APPROCHE SYSTÉMIQUE
    # =========================================================================
    s9 = prs.slides.add_slide(blank_layout)
    draw_bg(s9, BLANC)
    apply_transition(s9, "fade")
    add_official_header(s9, "NOTRE APPROCHE")

    approches = [
        ("FORMER", "Renforcer les compétences des acteurs territoriaux.", VERT_TERRITOIRE),
        ("ACCOMPAGNER", "Soutenir les acteurs à chaque étape de leur développement.", OR_RAYONNEMENT),
        ("INNOVER", "Proposer des solutions adaptées et durables aux enjeux locaux.", BLEU_INSTITUTIONNEL),
        ("TRANSFORMER LES TERRITOIRES", "Créer un impact positif et durable sur les communautés.", VERT_FORET)
    ]
    for i, (app_titre, app_desc, app_col) in enumerate(approches):
        x = Inches(0.8 + i * 2.95)
        y = Inches(2.3)
        ac = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.7), Inches(3.9))
        ac.fill.solid()
        ac.fill.fore_color.rgb = GRIS_FOND
        ac.line.color.rgb = app_col
        ac.line.width = Pt(1.5)
        
        actf = ac.text_frame
        actf.word_wrap = True
        actf.margin_top = Inches(0.4)
        actf.margin_left = Inches(0.2)
        actf.margin_right = Inches(0.2)

        p1 = actf.paragraphs[0]
        p1.text = f"0{i+1}"
        p1.font.bold = True
        p1.font.size = Pt(24)
        p1.font.color.rgb = app_col
        p1.alignment = PP_ALIGN.CENTER

        p2 = actf.add_paragraph()
        p2.text = f"\n{app_titre}"
        p2.font.bold = True
        p2.font.size = Pt(13)
        p2.font.color.rgb = app_col
        p2.alignment = PP_ALIGN.CENTER

        p3 = actf.add_paragraph()
        p3.text = f"\n\n{app_desc}"
        p3.font.size = Pt(11)
        p3.font.color.rgb = NOIR_TEXTE
        p3.alignment = PP_ALIGN.CENTER

    add_bottom_bar(s9)

    # =========================================================================
    # SLIDE 10 : NOS SERVICES
    # =========================================================================
    s10 = prs.slides.add_slide(blank_layout)
    draw_bg(s10, BLANC)
    apply_transition(s10, "push")
    add_official_header(s10, "NOS SERVICES")

    services = [
        ("FORMATION", "Sessions adaptées aux besoins des acteurs.", BLEU_INSTITUTIONNEL),
        ("CONSEIL", "Accompagnement stratégique et managérial.", VERT_TERRITOIRE),
        ("RECHERCHE", "Études et analyses territoriales approfondies.", OR_RAYONNEMENT),
        ("ACCOMPAGNEMENT", "Suivi et appui personnalisé sur le terrain.", VERT_TERRITOIRE),
        ("ÉTUDES & ÉVALUATIONS", "Diagnostics et évaluations d'impact.", OR_RAYONNEMENT),
        ("RENFORCEMENT DES CAPACITÉS", "Développement continu des compétences.", BLEU_INSTITUTIONNEL)
    ]
    for idx, (s_title, s_desc, s_col) in enumerate(services):
        row = idx // 3
        col = idx % 3
        x = Inches(0.8 + col * 4.0)
        y = Inches(2.3 + row * 1.9)
        sc = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.7), Inches(1.6))
        sc.fill.solid()
        sc.fill.fore_color.rgb = GRIS_FOND
        sc.line.color.rgb = s_col
        sc.line.width = Pt(1.5)

        sctf = sc.text_frame
        sctf.word_wrap = True
        sctf.margin_top = Inches(0.2)
        sctf.margin_left = Inches(0.25)

        p1 = sctf.paragraphs[0]
        p1.text = f"✦  {s_title}"
        p1.font.bold = True
        p1.font.size = Pt(13)
        p1.font.color.rgb = s_col

        p2 = sctf.add_paragraph()
        p2.text = s_desc
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = NOIR_TEXTE

    add_bottom_bar(s10)

    # =========================================================================
    # SLIDE 11 : NOTRE MÉTHODE D'INTERVENTION (5 ÉTAPES)
    # =========================================================================
    s11 = prs.slides.add_slide(blank_layout)
    draw_bg(s11, BLANC)
    apply_transition(s11, "wipe")
    add_official_header(s11, "NOTRE MÉTHODE D'INTERVENTION")

    method_steps = [
        ("01", "DIAGNOSTIC", "Analyser les besoins et les enjeux locaux."),
        ("02", "CO-CONSTRUCTION", "Définir ensemble les solutions adaptées."),
        ("03", "MISE EN ŒUVRE", "Mettre en œuvre les actions sur le terrain."),
        ("04", "SUIVI & ÉVALUATION", "Assurer un suivi rigoureux et participatif."),
        ("05", "CAPITALISATION", "Valoriser les résultats et les expériences.")
    ]
    for idx, (num, step_title, step_desc) in enumerate(method_steps):
        x = Inches(0.8 + idx * 2.35)
        y = Inches(2.4)
        mc = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.2), Inches(3.8))
        mc.fill.solid()
        mc.fill.fore_color.rgb = GRIS_FOND
        mc.line.color.rgb = BLEU_INSTITUTIONNEL if idx % 2 == 0 else VERT_TERRITOIRE
        mc.line.width = Pt(1.5)

        mctf = mc.text_frame
        mctf.word_wrap = True
        mctf.margin_top = Inches(0.3)
        mctf.margin_left = Inches(0.15)
        mctf.margin_right = Inches(0.15)

        p1 = mctf.paragraphs[0]
        p1.text = num
        p1.font.bold = True
        p1.font.size = Pt(22)
        p1.font.color.rgb = OR_RAYONNEMENT
        p1.alignment = PP_ALIGN.CENTER

        p2 = mctf.add_paragraph()
        p2.text = f"\n{step_title}"
        p2.font.bold = True
        p2.font.size = Pt(11)
        p2.font.color.rgb = BLEU_INSTITUTIONNEL
        p2.alignment = PP_ALIGN.CENTER

        p3 = mctf.add_paragraph()
        p3.text = f"\n\n{step_desc}"
        p3.font.size = Pt(10)
        p3.font.color.rgb = NOIR_TEXTE
        p3.alignment = PP_ALIGN.CENTER

    add_bottom_bar(s11)

    # =========================================================================
    # SLIDE 12 : ESPACES MUTUALISÉS & CHAMBRES D'HÔTES (THIÈS & SALY)
    # =========================================================================
    s12 = prs.slides.add_slide(blank_layout)
    draw_bg(s12, BLANC)
    apply_transition(s12, "fade")
    add_official_header(s12, "ESPACES MUTUALISÉS & CHAMBRES D'HÔTES")

    # Left: Salles
    if os.path.exists(img_seminaire):
        s12.shapes.add_picture(img_seminaire, Inches(0.8), Inches(2.2), width=Inches(5.6))
    c_salle = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.7), Inches(5.6), Inches(1.8))
    c_salle.fill.solid()
    c_salle.fill.fore_color.rgb = GRIS_FOND
    c_salle.line.color.rgb = VERT_TERRITOIRE
    stf = c_salle.text_frame
    stf.word_wrap = True
    sp1 = stf.paragraphs[0]
    sp1.text = "🏛️ Salles de Séminaires & Réunions (Thiès ZAC Nord)"
    sp1.font.bold = True
    sp1.font.size = Pt(12)
    sp1.font.color.rgb = VERT_FORET
    sp2 = stf.add_paragraph()
    sp2.text = "• Climatisation, Vidéoprojecteur, Sonorisation et Wifi très haut débit.\n• Modulable pour Séminaires, Réunions, Assemblées Générales et Formations."
    sp2.font.size = Pt(10)
    sp2.font.color.rgb = NOIR_TEXTE

    # Right: Chambres
    if os.path.exists(img_chambre):
        s12.shapes.add_picture(img_chambre, Inches(6.8), Inches(2.2), width=Inches(5.7))
    c_chambre = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(4.7), Inches(5.7), Inches(1.8))
    c_chambre.fill.solid()
    c_chambre.fill.fore_color.rgb = GRIS_FOND
    c_chambre.line.color.rgb = OR_RAYONNEMENT
    htf = c_chambre.text_frame
    htf.word_wrap = True
    hp1 = htf.paragraphs[0]
    hp1.text = "🏡 Chambres d'Hôtes (Thiès & Saly Portudal)"
    hp1.font.bold = True
    hp1.font.size = Pt(12)
    hp1.font.color.rgb = OR_RAYONNEMENT
    hp2 = htf.add_paragraph()
    hp2.text = "• Résidence Thiès (missions professionnelles) et Villa Saly (séjours & retraites).\n• Petit-déjeuner bio local inclus, navette Aéroport International AIBD."
    hp2.font.size = Pt(10)
    hp2.font.color.rgb = NOIR_TEXTE

    # =========================================================================
    # SLIDE 13 : APPEL À PARTENARIAT & ACTEURS
    # =========================================================================
    s13 = prs.slides.add_slide(blank_layout)
    draw_bg(s13, BLANC)
    apply_transition(s13, "push")
    add_official_header(s13, "APPEL À PARTENARIAT")

    # Grand Cadre Central
    p_box = s13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.2), Inches(11.7), Inches(1.5))
    p_box.fill.solid()
    p_box.fill.fore_color.rgb = VERT_PALE
    p_box.line.color.rgb = VERT_TERRITOIRE
    ptf = p_box.text_frame
    ptf.word_wrap = True
    ptf.margin_left = Inches(0.4)
    ptf.margin_top = Inches(0.2)

    pp1 = ptf.paragraphs[0]
    pp1.text = "Construisons ensemble des territoires durables."
    pp1.font.bold = True
    pp1.font.size = Pt(18)
    pp1.font.color.rgb = VERT_FORET

    pp2 = ptf.add_paragraph()
    pp2.text = "APATAM@E est ouverte à toutes formes de partenariats avec les institutions, les collectivités territoriales, les organisations de la société civile, le secteur privé et les partenaires techniques et financiers."
    pp2.font.size = Pt(12)
    pp2.font.color.rgb = NOIR_TEXTE

    # 5 Catégories de Partenaires
    partenaires_cats = [
        ("Institutions publiques", "Ministères & Agences"),
        ("Collectivités territoriales", "Mairies & Départements"),
        ("Partenaires techniques", "Bailleurs & ONG"),
        ("Secteur privé", "Entreprises engagées"),
        ("Société civile", "Associations & Jeunesse")
    ]
    for idx, (cat_t, cat_sub) in enumerate(partenaires_cats):
        x = Inches(0.8 + idx * 2.35)
        y = Inches(4.0)
        c_part = s13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.2), Inches(2.2))
        c_part.fill.solid()
        c_part.fill.fore_color.rgb = GRIS_FOND
        c_part.line.color.rgb = BORDER_GRIS

        cptf = c_part.text_frame
        cptf.word_wrap = True
        cptf.margin_top = Inches(0.4)
        cptf.margin_left = Inches(0.15)
        cptf.margin_right = Inches(0.15)

        p1 = cptf.paragraphs[0]
        p1.text = f"✦\n{cat_t}"
        p1.font.bold = True
        p1.font.size = Pt(12)
        p1.font.color.rgb = BLEU_INSTITUTIONNEL
        p1.alignment = PP_ALIGN.CENTER

        p2 = cptf.add_paragraph()
        p2.text = f"\n{cat_sub}"
        p2.font.size = Pt(9.5)
        p2.font.color.rgb = GRIS_TEXTE
        p2.alignment = PP_ALIGN.CENTER

    add_bottom_bar(s13)

    # =========================================================================
    # SLIDE 14 : CONTACTS & COORDONNÉES OFFICIELLES
    # =========================================================================
    s14 = prs.slides.add_slide(blank_layout)
    draw_bg(s14, BLANC)
    apply_transition(s14, "fade")
    add_official_header(s14, "CONTACTS & COORDONNÉES")

    # Left Info Card
    c_contact = s14.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.2), Inches(7.5), Inches(4.2))
    c_contact.fill.solid()
    c_contact.fill.fore_color.rgb = GRIS_FOND
    c_contact.line.color.rgb = VERT_TERRITOIRE
    c_contact.line.width = Pt(1.5)

    ctf = c_contact.text_frame
    ctf.word_wrap = True
    ctf.margin_left = Inches(0.4)
    ctf.margin_top = Inches(0.3)

    cp1 = ctf.paragraphs[0]
    cp1.text = "APATAM@E — Direction Générale"
    cp1.font.bold = True
    cp1.font.size = Pt(16)
    cp1.font.color.rgb = BLEU_INSTITUTIONNEL

    cp2 = ctf.add_paragraph()
    cp2.text = "Présidente Fondatrice : Adama Mbengue\n"
    cp2.font.size = Pt(12)
    cp2.font.bold = True
    cp2.font.color.rgb = VERT_FORET

    cp3 = ctf.add_paragraph()
    cp3.text = "📍 Siège social : ZAC NORD Thiès Nº2688, Rond Point Mosquée, Thiès - Sénégal\n"
    cp3.font.size = Pt(11.5)
    cp3.font.color.rgb = NOIR_TEXTE

    cp4 = ctf.add_paragraph()
    cp4.text = "📞 Bureau : +221 33 999 28 52  |  Portable : +221 77 510 20 38\n"
    cp4.font.size = Pt(11.5)
    cp4.font.color.rgb = NOIR_TEXTE

    cp5 = ctf.add_paragraph()
    cp5.text = "📧 E-mails : contact@apatame.com  |  secretariat@apatame.org\n"
    cp5.font.size = Pt(11.5)
    cp5.font.color.rgb = NOIR_TEXTE

    cp6 = ctf.add_paragraph()
    cp6.text = "🌐 Site Web : www.apatame.org (ou http://localhost:5173)\n🏡 Espaces & Chambres d'Hôtes : Thiès & Saly Portudal"
    cp6.font.size = Pt(11.5)
    cp6.font.bold = True
    cp6.font.color.rgb = OR_RAYONNEMENT

    # Right: Carte Scan QR & Présidente
    if os.path.exists(img_pres):
        s14.shapes.add_picture(img_pres, Inches(8.7), Inches(2.2), width=Inches(3.8))

    add_bottom_bar(s14)

    # =========================================================================
    # SLIDE 15 : MERCI POUR VOTRE ATTENTION (CLOSING)
    # =========================================================================
    s15 = prs.slides.add_slide(blank_layout)
    draw_bg(s15, BLANC)
    apply_transition(s15, "fade")

    if os.path.exists(img_merci):
        s15.shapes.add_picture(img_merci, 0, 0, width=Inches(13.333), height=Inches(7.5))
    else:
        # Fallback élégant
        draw_bg(s15, VERT_FORET)
        if os.path.exists(logo_path):
            s15.shapes.add_picture(logo_path, Inches(5.8), Inches(1.5), height=Inches(1.5))
        
        tb_m = s15.shapes.add_textbox(Inches(2.0), Inches(3.2), Inches(9.33), Inches(2.5))
        tf_m = tb_m.text_frame
        p = tf_m.paragraphs[0]
        p.text = "MERCI"
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = OR_CLAIR
        p.alignment = PP_ALIGN.CENTER

        p2 = tf_m.add_paragraph()
        p2.text = "POUR VOTRE ATTENTION"
        p2.font.size = Pt(20)
        p2.font.bold = True
        p2.font.color.rgb = BLANC
        p2.alignment = PP_ALIGN.CENTER

        p3 = tf_m.add_paragraph()
        p3.text = "\nEnsemble, bâtissons des territoires durables et prospères."
        p3.font.size = Pt(14)
        p3.font.italic = True
        p3.font.color.rgb = BLANC
        p3.alignment = PP_ALIGN.CENTER

    # Save to all destination paths
    filename = "APATAM_E_Presentation_Officielle_Complete_2026.pptx"
    
    # 1. Parent folder
    parent_dir = os.path.abspath(os.path.join(base_dir, ".."))
    parent_path = os.path.join(parent_dir, filename)
    prs.save(parent_path)
    print(f"[OK] Saved to parent directory: {parent_path}")

    # 2. Downloads folder
    user_home = os.path.expanduser("~")
    downloads_path = os.path.join(user_home, "Downloads", filename)
    prs.save(downloads_path)
    print(f"[OK] Saved to Downloads: {downloads_path}")

    # 3. Project folder
    project_path = os.path.join(base_dir, filename)
    prs.save(project_path)
    print(f"[OK] Saved to Project: {project_path}")

    # 4. E:\noppal\apatamae folder if available
    if os.path.exists(e_drive_dir):
        e_path = os.path.join(e_drive_dir, filename)
        prs.save(e_path)
        print(f"[OK] Saved to E:\\noppal\\apatamae: {e_path}")

if __name__ == "__main__":
    create_master_presentation()
