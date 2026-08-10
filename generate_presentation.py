import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # 16:9 Widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6] # Blank slide

    # Color Palette
    DARK_GREEN = RGBColor(1, 45, 29)       # #012D1D
    MID_GREEN  = RGBColor(27, 67, 50)      # #1B4332
    NAVY_BLUE  = RGBColor(57, 96, 147)     # #396093
    GOLD       = RGBColor(218, 165, 32)    # #DAA520
    LIGHT_BG   = RGBColor(248, 249, 250)   # #F8F9FA
    WHITE      = RGBColor(255, 255, 255)
    DARK_TEXT  = RGBColor(25, 28, 29)
    GRAY_TEXT  = RGBColor(100, 110, 105)
    CARD_BG    = RGBColor(255, 255, 255)
    MUTED_GREEN= RGBColor(165, 208, 185)

    base_dir = os.path.dirname(os.path.abspath(__file__))
    images_dir = os.path.join(base_dir, "public", "images")
    logo_path = os.path.join(base_dir, "public", "logo.png")

    def add_bg(slide, color=LIGHT_BG):
        # Full slide background shape
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = color
        bg.line.fill.background() # No border
        return bg

    def add_header(slide, title, category="APATAM@E • INSTITUTIONAL", dark_mode=False):
        # Top banner category
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.35))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = category.upper()
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = GOLD if dark_mode else NAVY_BLUE
        
        # Main Title
        tb2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(10), Inches(0.8))
        tf2 = tb2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = title
        p2.font.size = Pt(24)
        p2.font.bold = True
        p2.font.color.rgb = WHITE if dark_mode else DARK_GREEN

        # Optional Top-Right Mini Logo
        if os.path.exists(logo_path):
            slide.shapes.add_picture(logo_path, Inches(11.8), Inches(0.4), width=Inches(0.8))

    # =========================================================================
    # SLIDE 1: COVER / TITRE
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_slide_layout)
    add_bg(slide1, DARK_GREEN)

    # Decorative abstract background map if available
    abstract_img = os.path.join(images_dir, "a_modern_glowing_abstract_3d_map_of_senegal_composed_of_floating_dots_and_nodes_screen.png")
    if os.path.exists(abstract_img):
        pic = slide1.shapes.add_picture(abstract_img, Inches(7.0), Inches(0.8), width=Inches(6.0))

    # Logo on Title Slide
    if os.path.exists(logo_path):
        slide1.shapes.add_picture(logo_path, Inches(1.0), Inches(1.0), width=Inches(1.6))

    # Category Pill
    pill = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(2.9), Inches(3.8), Inches(0.45))
    pill.fill.solid()
    pill.fill.fore_color.rgb = MID_GREEN
    pill.line.color.rgb = GOLD
    p = pill.text_frame.paragraphs[0]
    p.text = "PLATEFORME INSTITUTIONNELLE & DÉVELOPPEMENT"
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER

    # Main Title
    t_box = slide1.shapes.add_textbox(Inches(0.9), Inches(3.5), Inches(7.5), Inches(2.2))
    tf = t_box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = "APATAM@E"
    p1.font.size = Pt(44)
    p1.font.bold = True
    p1.font.color.rgb = GOLD

    p2 = tf.add_paragraph()
    p2.text = "Gouvernance, Territoires & Avenir"
    p2.font.size = Pt(28)
    p2.font.bold = True
    p2.font.color.rgb = WHITE

    p3 = tf.add_paragraph()
    p3.text = "Bâtir des territoires inclusifs, résilients et connectés aux savoirs."
    p3.font.size = Pt(13)
    p3.font.color.rgb = MUTED_GREEN

    # Footer banner on slide 1
    f_box = slide1.shapes.add_textbox(Inches(0.9), Inches(6.2), Inches(11.5), Inches(0.6))
    ftf = f_box.text_frame
    fp = ftf.paragraphs[0]
    fp.text = "Présidente Fondatrice : Adama Mbengue  •  Siège : ZAC NORD Thiès Nº2688, Sénégal  •  contact@apatame.com"
    fp.font.size = Pt(10)
    fp.font.color.rgb = WHITE

    # =========================================================================
    # SLIDE 2: QUI SOMMES-NOUS & NOTRE ADN
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_slide_layout)
    add_bg(slide2, LIGHT_BG)
    add_header(slide2, "Notre Identité & Notre Vision Stratégique", "NOTRE ADN & CONVICTIONS")

    # Left Image
    dna_img = os.path.join(images_dir, "high_quality_artistic_photography_of_senegalese_nature_and_local_citizens_screen.png")
    if os.path.exists(dna_img):
        slide2.shapes.add_picture(dna_img, Inches(0.8), Inches(1.8), width=Inches(4.5))

    # Right Content Cards
    cards_data = [
        ("Notre Mission Fondatrice", "Former, informer, accompagner et mobiliser les citoyens et les acteurs territoriaux pour une gouvernance responsable et un développement local inclusif et durable au Sénégal."),
        ("Notre Approche Endogène", "Allier la richesse des traditions positives sénégalaises aux outils modernes de gouvernance, du digital et de la transition écologique."),
        ("Notre Rayonnement Territorial", "Un ancrage fort à Thiès (Siège national ZAC Nord), à Dakar, à Saly ainsi que dans les terroirs du Nord, du Fleuve, du Centre et du Sud.")
    ]

    top_pos = 1.8
    for title, desc in cards_data:
        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.6), Inches(top_pos), Inches(6.9), Inches(1.5))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RGBColor(220, 225, 220)
        
        ctf = card.text_frame
        ctf.word_wrap = True
        ctf.margin_left = Inches(0.3)
        ctf.margin_top = Inches(0.2)
        
        cp1 = ctf.paragraphs[0]
        cp1.text = title
        cp1.font.bold = True
        cp1.font.size = Pt(14)
        cp1.font.color.rgb = DARK_GREEN
        
        cp2 = ctf.add_paragraph()
        cp2.text = desc
        cp2.font.size = Pt(10.5)
        cp2.font.color.rgb = DARK_TEXT
        
        top_pos += 1.7

    # =========================================================================
    # SLIDE 3: MOT DE LA PRÉSIDENTE FONDATRICE
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_slide_layout)
    add_bg(slide3, LIGHT_BG)
    add_header(slide3, "Le Mot de la Présidente Fondatrice", "GOUVERNANCE & INSPIRATION")

    # President Photo
    pres_img = os.path.join(images_dir, "professional_corporate_portrait_of_a_west_african_woman_leader_elegant_and_screen.png")
    if os.path.exists(pres_img):
        slide3.shapes.add_picture(pres_img, Inches(0.8), Inches(1.8), width=Inches(3.8))

    # Quote Box
    q_card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.9), Inches(1.8), Inches(7.6), Inches(4.8))
    q_card.fill.solid()
    q_card.fill.fore_color.rgb = DARK_GREEN
    q_card.line.color.rgb = GOLD
    q_card.line.width = Pt(1.5)

    qtf = q_card.text_frame
    qtf.word_wrap = True
    qtf.margin_left = Inches(0.5)
    qtf.margin_right = Inches(0.5)
    qtf.margin_top = Inches(0.5)

    qp1 = qtf.paragraphs[0]
    qp1.text = "« Notre mission est de promouvoir des modèles endogènes de gouvernance et de développement, en réconciliant modernité et traditions positives pour une Afrique debout. »"
    qp1.font.size = Pt(17)
    qp1.font.bold = True
    qp1.font.italic = True
    qp1.font.color.rgb = WHITE

    qp2 = qtf.add_paragraph()
    qp2.text = "\nChaque citoyen, ancré dans ses valeurs, détient le potentiel d'être le moteur d'une transformation collective durable. À travers APATAM@E, nous construisons des passerelles entre la recherche, les institutions et les communautés de base."
    qp2.font.size = Pt(12)
    qp2.font.color.rgb = MUTED_GREEN

    qp3 = qtf.add_paragraph()
    qp3.text = "\n\nAdama Mbengue"
    qp3.font.size = Pt(16)
    qp3.font.bold = True
    qp3.font.color.rgb = GOLD

    qp4 = qtf.add_paragraph()
    qp4.text = "Présidente Fondatrice — APATAM@E Sénégal"
    qp4.font.size = Pt(10)
    qp4.font.color.rgb = WHITE

    # =========================================================================
    # SLIDE 4: NOS 7 PRINCIPES D'ACTION & VALEURS
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_slide_layout)
    add_bg(slide4, LIGHT_BG)
    add_header(slide4, "Nos Principes d'Action & Valeurs Fondatrices", "FONDEMENTS ÉTHIQUES")

    values = [
        ("L'Équité", "Justice sociale, égalité des chances et inclusion de genre."),
        ("L'Intégrité", "Transparence, redevabilité et exemplarité de gestion."),
        ("L'Inclusion", "Participation active de tous les terroirs et des jeunes."),
        ("L'Innovation", "Savoirs endogènes couplés aux technologies modernes."),
        ("La Collaboration", "Synergies partenariales, co-construction citoyenne."),
        ("Paix & Sécurité", "Médiation communautaire, cohésion et vivre-ensemble."),
        ("Transmission", "Dialogue permanent entre aînés et générations montantes.")
    ]

    # Render 7 Cards across 2 rows (4 on top, 3 on bottom)
    card_w = Inches(2.7)
    card_h = Inches(2.2)
    
    # Row 1 (4 items)
    for i in range(4):
        x = Inches(0.8 + i * 2.95)
        y = Inches(1.8)
        c = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h)
        c.fill.solid()
        c.fill.fore_color.rgb = WHITE
        c.line.color.rgb = RGBColor(220, 225, 220)
        ctf = c.text_frame
        ctf.word_wrap = True
        ctf.margin_top = Inches(0.25)
        p1 = ctf.paragraphs[0]
        p1.text = f"0{i+1}. {values[i][0]}"
        p1.font.bold = True
        p1.font.size = Pt(13)
        p1.font.color.rgb = DARK_GREEN
        p2 = ctf.add_paragraph()
        p2.text = f"\n{values[i][1]}"
        p2.font.size = Pt(10)
        p2.font.color.rgb = DARK_TEXT

    # Row 2 (3 items)
    for j in range(3):
        idx = 4 + j
        x = Inches(1.5 + j * 3.6)
        y = Inches(4.3)
        c = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.3), card_h)
        c.fill.solid()
        c.fill.fore_color.rgb = WHITE
        c.line.color.rgb = RGBColor(220, 225, 220)
        ctf = c.text_frame
        ctf.word_wrap = True
        ctf.margin_top = Inches(0.25)
        p1 = ctf.paragraphs[0]
        p1.text = f"0{idx+1}. {values[idx][0]}"
        p1.font.bold = True
        p1.font.size = Pt(13)
        p1.font.color.rgb = DARK_GREEN if j != 1 else NAVY_BLUE
        p2 = ctf.add_paragraph()
        p2.text = f"\n{values[idx][1]}"
        p2.font.size = Pt(10)
        p2.font.color.rgb = DARK_TEXT

    # =========================================================================
    # SLIDE 5: NOS 6 PROGRAMMES MAJEURS
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_slide_layout)
    add_bg(slide5, LIGHT_BG)
    add_header(slide5, "Nos 6 Programmes & Piliers Stratégiques", "IMPACT & TRANSFORMATION")

    programs = [
        ("01. Éducation & STEM", "Réduction des fractures éducatives, bourses d'excellence et mentorat scientifique féminin."),
        ("02. Académie du Leadership", "Formation de la nouvelle génération de leaders territoriaux et décideurs engagés."),
        ("03. Résilience Agricole", "Agroécologie sahélienne, irrigation solaire et valorisation des filières durables."),
        ("04. Observatoire Public", "Recherche-action, baromètre citoyen et plaidoyer éclairé pour les politiques locales."),
        ("05. Paix & Sécurité", "Comités locaux de médiation, prévention des conflits et cohésion territoriale."),
        ("06. Échanges Aînés & Jeunes", "Sessions sous l'Apatam, transmission des savoirs endogènes et mentorat moral.")
    ]

    p_w = Inches(3.7)
    p_h = Inches(2.2)
    for idx, (p_title, p_desc) in enumerate(programs):
        row = idx // 3
        col = idx % 3
        x = Inches(0.8 + col * 4.0)
        y = Inches(1.8 + row * 2.5)

        card = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, p_w, p_h)
        card.fill.solid()
        # Alternate subtly between dark green and white
        if idx in [0, 4]:
            card.fill.fore_color.rgb = DARK_GREEN
            card.line.color.rgb = DARK_GREEN
            title_col = GOLD
            desc_col = WHITE
        elif idx in [2, 5]:
            card.fill.fore_color.rgb = WHITE
            card.line.color.rgb = MID_GREEN
            title_col = DARK_GREEN
            desc_col = DARK_TEXT
        else:
            card.fill.fore_color.rgb = WHITE
            card.line.color.rgb = NAVY_BLUE
            title_col = NAVY_BLUE
            desc_col = DARK_TEXT

        ctf = card.text_frame
        ctf.word_wrap = True
        ctf.margin_top = Inches(0.2)
        ctf.margin_left = Inches(0.25)
        ctf.margin_right = Inches(0.25)

        p1 = ctf.paragraphs[0]
        p1.text = p_title
        p1.font.bold = True
        p1.font.size = Pt(13)
        p1.font.color.rgb = title_col

        p2 = ctf.add_paragraph()
        p2.text = f"\n{p_desc}"
        p2.font.size = Pt(10)
        p2.font.color.rgb = desc_col

    # =========================================================================
    # SLIDE 6: GOUVERNANCE & ORGANIGRAMME
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_slide_layout)
    add_bg(slide6, LIGHT_BG)
    add_header(slide6, "Structure Organisationnelle & Gouvernance", "GOUVERNANCE INTERNE")

    # Level 1: Conseil d'Orientation Stratégique
    c1 = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), Inches(1.6), Inches(4.3), Inches(0.85))
    c1.fill.solid()
    c1.fill.fore_color.rgb = NAVY_BLUE
    c1.line.color.rgb = NAVY_BLUE
    p = c1.text_frame.paragraphs[0]
    p.text = "Conseil d'Orientation Stratégique"
    p.font.bold = True
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Level 2: Direction Générale & Secrétariat
    c2 = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), Inches(2.7), Inches(4.3), Inches(0.85))
    c2.fill.solid()
    c2.fill.fore_color.rgb = DARK_GREEN
    c2.line.color.rgb = GOLD
    p = c2.text_frame.paragraphs[0]
    p.text = "Direction Générale (Adama Mbengue)"
    p.font.bold = True
    p.font.size = Pt(12)
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER

    c2_sec = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.1), Inches(2.7), Inches(3.4), Inches(0.85))
    c2_sec.fill.solid()
    c2_sec.fill.fore_color.rgb = WHITE
    c2_sec.line.color.rgb = RGBColor(200, 210, 200)
    p = c2_sec.text_frame.paragraphs[0]
    p.text = "Secrétariat de Direction\n(Agenda, courrier & accueil)"
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = DARK_GREEN
    p.alignment = PP_ALIGN.CENTER

    # Level 3: 4 Directions Opérationnelles
    dirs = [
        "Dir. des Programmes\n& Formations",
        "Dir. de la Recherche\n& Observatoire",
        "Dir. des Partenariats\n& Communication",
        "Dir. Administrative\n& Financière"
    ]
    for i, d in enumerate(dirs):
        cd = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + i * 2.95), Inches(3.8), Inches(2.7), Inches(1.1))
        cd.fill.solid()
        cd.fill.fore_color.rgb = MID_GREEN
        cd.line.color.rgb = MID_GREEN
        p = cd.text_frame.paragraphs[0]
        p.text = d
        p.font.bold = True
        p.font.size = Pt(10.5)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # Level 4: 5 Services d'Appui
    supp = ["Comptabilité & Finances", "Ressources Humaines", "Informatique & Digital", "Logistique & Patrimoine", "Suivi-Éval & Qualité"]
    for j, s in enumerate(supp):
        cs = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + j * 2.36), Inches(5.2), Inches(2.2), Inches(0.8))
        cs.fill.solid()
        cs.fill.fore_color.rgb = WHITE
        cs.line.color.rgb = RGBColor(200, 210, 200)
        p = cs.text_frame.paragraphs[0]
        p.text = s
        p.font.bold = True
        p.font.size = Pt(9)
        p.font.color.rgb = DARK_GREEN
        p.alignment = PP_ALIGN.CENTER

    # Bottom Mission Bar
    m_bar = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(6.3), Inches(11.73), Inches(0.6))
    m_bar.fill.solid()
    m_bar.fill.fore_color.rgb = DARK_GREEN
    m_bar.line.fill.background()
    p = m_bar.text_frame.paragraphs[0]
    p.text = "MISSION : Former, informer, accompagner et mobiliser les citoyens pour un développement inclusif et durable."
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # =========================================================================
    # SLIDE 7: ESPACES MUTUALISÉS & CHAMBRES D'HÔTES
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_slide_layout)
    add_bg(slide7, LIGHT_BG)
    add_header(slide7, "Espaces Mutualisés & Chambres d'Hôtes", "ACCUEIL & INFRASTRUCTURES")

    # Left: Salles Équipées
    s_card = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    s_card.fill.solid()
    s_card.fill.fore_color.rgb = WHITE
    s_card.line.color.rgb = DARK_GREEN
    s_card.line.width = Pt(1.5)
    stf = s_card.text_frame
    stf.word_wrap = True
    stf.margin_left = Inches(0.4)
    stf.margin_top = Inches(0.3)

    sp1 = stf.paragraphs[0]
    sp1.text = "🏛️ Salles de Séminaires & Réunions"
    sp1.font.bold = True
    sp1.font.size = Pt(16)
    sp1.font.color.rgb = DARK_GREEN

    sp2 = stf.add_paragraph()
    sp2.text = "\n• Localisation : Maison de la Citoyenneté (ZAC Nord Thiès)\n• Types d'activités : Séminaires, Réunions, Assemblées Générales, Formations & Ateliers de travail.\n• Équipements inclus : Climatisation, Vidéoprojecteur interactif, Sonorisation haute fidélité, Wifi haut débit et pause-café locale.\n• Modularité : Configuration selon l'effectif exact (jusqu'à 150 personnes)."
    sp2.font.size = Pt(11)
    sp2.font.color.rgb = DARK_TEXT

    # Right: Chambres d'Hôtes
    h_card = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8))
    h_card.fill.solid()
    h_card.fill.fore_color.rgb = DARK_GREEN
    h_card.line.color.rgb = GOLD
    h_card.line.width = Pt(1.5)
    htf = h_card.text_frame
    htf.word_wrap = True
    htf.margin_left = Inches(0.4)
    htf.margin_top = Inches(0.3)

    hp1 = htf.paragraphs[0]
    hp1.text = "🏡 Chambres d'Hôtes (Thiès & Saly)"
    hp1.font.bold = True
    hp1.font.size = Pt(16)
    hp1.font.color.rgb = GOLD

    hp2 = htf.add_paragraph()
    hp2.text = "\n• Destination Thiès : Résidence ZAC Nord Nº2688 (calme, sérénité, idéal pour missions professionnelles et universitaires).\n• Destination Saly Portudal : Villa d'hôtes balnéaire, piscine et cadre reposant pour retraites de travail et séjours détente.\n• Formules : Chambres Individuelles Confort, Suites Master et Privatisation complète de villa.\n• Services : Petit-déjeuner bio local inclus, navette Aéroport International AIBD."
    hp2.font.size = Pt(11)
    hp2.font.color.rgb = WHITE

    # =========================================================================
    # SLIDE 8: IMPACT & RÉSULTATS CLÉS
    # =========================================================================
    slide8 = prs.slides.add_slide(blank_slide_layout)
    add_bg(slide8, DARK_GREEN)
    add_header(slide8, "Mesure de la Performance & Transformation", "IMPACT & RÉSULTATS", dark_mode=True)

    metrics_data = [
        ("25,000+", "Bénéficiaires Directs", "Agricultrices, lycéennes, jeunes leaders et élus locaux mobilisés."),
        ("48", "Périmètres Décarbonés", "Systèmes d'irrigation solaire et fermes agroécologiques autonomes."),
        ("94%", "Taux d'Autonomie", "Des coopératives féminines accompagnées pérennes après 24 mois."),
        ("14", "Notes de Recherche", "Publications scientifiques et baromètres de gouvernance en open data.")
    ]

    for idx, (m_val, m_lbl, m_desc) in enumerate(metrics_data):
        x = Inches(0.8 + idx * 2.95)
        y = Inches(2.2)
        m_card = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.7), Inches(3.8))
        m_card.fill.solid()
        m_card.fill.fore_color.rgb = MID_GREEN
        m_card.line.color.rgb = GOLD if idx == 0 else RGBColor(50, 90, 70)
        
        mtf = m_card.text_frame
        mtf.word_wrap = True
        mtf.margin_top = Inches(0.3)
        mtf.margin_left = Inches(0.2)
        mtf.margin_right = Inches(0.2)

        p1 = mtf.paragraphs[0]
        p1.text = m_val
        p1.font.bold = True
        p1.font.size = Pt(28)
        p1.font.color.rgb = GOLD
        p1.alignment = PP_ALIGN.CENTER

        p2 = mtf.add_paragraph()
        p2.text = f"\n{m_lbl}"
        p2.font.bold = True
        p2.font.size = Pt(13)
        p2.font.color.rgb = WHITE
        p2.alignment = PP_ALIGN.CENTER

        p3 = mtf.add_paragraph()
        p3.text = f"\n{m_desc}"
        p3.font.size = Pt(9.5)
        p3.font.color.rgb = MUTED_GREEN
        p3.alignment = PP_ALIGN.CENTER

    # =========================================================================
    # SLIDE 9: ALLIANCES & PARTENARIATS STRATÉGIQUES
    # =========================================================================
    slide9 = prs.slides.add_slide(blank_slide_layout)
    add_bg(slide9, LIGHT_BG)
    add_header(slide9, "Partenaires & Alliances Stratégiques", "COOPÉRATION & SYNERGIES")

    p_categories = [
        ("🏛️ Ministères & Agences Publiques", [
            "Ministère de l'Environnement et du Développement Durable",
            "Ministère de l'Enseignement Supérieur et de la Recherche",
            "Agence Nationale des Énergies Renouvelables (ANER)",
            "Collectivités territoriales et Mairies partenaires"
        ]),
        ("🌍 Bailleurs & Coopération Internationale", [
            "Union Européenne (Programmes Climat & Jeunesse)",
            "PNUD Sénégal (Fonds pour l'Environnement Mondial)",
            "Coopération Allemande au Développement (GIZ)",
            "Fondations philanthropiques pour l'Afrique"
        ]),
        ("🎓 Universités & Pôles de Recherche", [
            "Université Cheikh Anta Diop (UCAD - Dakar)",
            "Université Gaston Berger (UGB - Saint-Louis)",
            "Institut Sénégalais de Recherches Agricoles (ISRA)",
            "Institut de Technologie Alimentaire (ITA)"
        ])
    ]

    for idx, (cat_title, p_list) in enumerate(p_categories):
        x = Inches(0.8 + idx * 3.95)
        y = Inches(1.8)
        p_card = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.7), Inches(4.8))
        p_card.fill.solid()
        p_card.fill.fore_color.rgb = WHITE
        p_card.line.color.rgb = RGBColor(220, 225, 220)
        
        ptf = p_card.text_frame
        ptf.word_wrap = True
        ptf.margin_top = Inches(0.3)
        ptf.margin_left = Inches(0.3)
        ptf.margin_right = Inches(0.3)

        p1 = ptf.paragraphs[0]
        p1.text = cat_title
        p1.font.bold = True
        p1.font.size = Pt(13)
        p1.font.color.rgb = DARK_GREEN

        for p_item in p_list:
            p_bullet = ptf.add_paragraph()
            p_bullet.text = f"• {p_item}"
            p_bullet.font.size = Pt(10.5)
            p_bullet.font.color.rgb = DARK_TEXT

    # =========================================================================
    # SLIDE 10: COORDONNÉES OFFICIELLES & CONTACT
    # =========================================================================
    slide10 = prs.slides.add_slide(blank_slide_layout)
    add_bg(slide10, DARK_GREEN)
    add_header(slide10, "Rejoindre l'Action & Nous Contacter", "CONTACT & RÉSERVATIONS", dark_mode=True)

    # Main Card
    main_c = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.8), Inches(10.33), Inches(4.8))
    main_c.fill.solid()
    main_c.fill.fore_color.rgb = MID_GREEN
    main_c.line.color.rgb = GOLD
    main_c.line.width = Pt(1.5)

    mtf = main_c.text_frame
    mtf.word_wrap = True
    mtf.margin_left = Inches(0.6)
    mtf.margin_right = Inches(0.6)
    mtf.margin_top = Inches(0.4)

    p1 = mtf.paragraphs[0]
    p1.text = "APATAM@E — Antigravity Institutional & Development"
    p1.font.bold = True
    p1.font.size = Pt(20)
    p1.font.color.rgb = GOLD

    p2 = mtf.add_paragraph()
    p2.text = "Présidente Fondatrice : Adama Mbengue"
    p2.font.size = Pt(13)
    p2.font.bold = True
    p2.font.color.rgb = WHITE

    p3 = mtf.add_paragraph()
    p3.text = "\n📍 Siège National : ZAC NORD Thiès Nº2688, Rond Point Mosquée, Thiès - Sénégal\n"
    p3.font.size = Pt(12)
    p3.font.color.rgb = MUTED_GREEN

    p4 = mtf.add_paragraph()
    p4.text = "📧 Courriels Officiels : contact@apatame.com  |  secretariat@apatame.org"
    p4.font.size = Pt(12)
    p4.font.color.rgb = WHITE

    p5 = mtf.add_paragraph()
    p5.text = "📞 Téléphones : Bureau : 33 999 28 52  |  Portable : 77 510 20 38"
    p5.font.size = Pt(12)
    p5.font.color.rgb = WHITE

    p6 = mtf.add_paragraph()
    p6.text = "🏡 Espaces & Chambres d'Hôtes : Thiès & Saly Portudal"
    p6.font.size = Pt(12)
    p6.font.color.rgb = GOLD

    p7 = mtf.add_paragraph()
    p7.text = "\n🌐 Plateforme Web : http://localhost:5173"
    p7.font.size = Pt(12)
    p7.font.bold = True
    p7.font.color.rgb = WHITE

    # Target save paths
    output_filename = "Presentation_APATAM_E_Institutionnelle_2026.pptx"
    
    # Path 1: Parent directory of workspace
    parent_dir = os.path.abspath(os.path.join(base_dir, ".."))
    parent_save_path = os.path.join(parent_dir, output_filename)
    prs.save(parent_save_path)
    print(f"[OK] Saved to parent directory: {parent_save_path}")

    # Path 2: Project directory
    local_save_path = os.path.join(base_dir, output_filename)
    prs.save(local_save_path)
    print(f"[OK] Saved to project directory: {local_save_path}")

    # Path 3: Downloads directory if accessible
    user_home = os.path.expanduser("~")
    downloads_dir = os.path.join(user_home, "Downloads")
    if os.path.exists(downloads_dir):
        downloads_save_path = os.path.join(downloads_dir, output_filename)
        prs.save(downloads_save_path)
        print(f"[OK] Saved to Downloads: {downloads_save_path}")

if __name__ == "__main__":
    create_presentation()
