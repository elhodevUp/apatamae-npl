import os
import sys
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml import parse_xml

def build_modern_animated_deck():
    prs = Presentation()
    # 16:9 Modern Widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette: Modern, Simple & Refined
    DARK_FOREST = RGBColor(1, 45, 29)      # #012D1D
    DEEP_GREEN  = RGBColor(15, 60, 42)     # #0F3C2A
    OCEAN_BLUE  = RGBColor(42, 82, 132)    # #2A5284
    ACCENT_GOLD = RGBColor(218, 165, 32)   # #DAA520
    BRIGHT_GOLD = RGBColor(245, 190, 44)   # #F5BE2C
    SOFT_BG     = RGBColor(250, 251, 252)  # #FAFBFC
    CARD_WHITE  = RGBColor(255, 255, 255)
    CARD_BORDER = RGBColor(228, 233, 230)
    DARK_TEXT   = RGBColor(18, 24, 22)
    BODY_TEXT   = RGBColor(70, 80, 75)
    MUTED_TEXT  = RGBColor(140, 150, 145)
    LIGHT_MINT  = RGBColor(180, 220, 200)

    base_dir = os.path.dirname(os.path.abspath(__file__))
    artifacts_dir = r"C:\Users\LENOVO E15\.gemini\antigravity-ide\brain\2584ad8b-cb21-4a7f-a2cd-17da14691d98"
    public_images = os.path.join(base_dir, "public", "images")
    logo_path = os.path.join(base_dir, "public", "logo.png")

    # Image Paths
    img_hero = os.path.join(artifacts_dir, "pptx_hero_senegal_1785866581210.png")
    img_seminaire = os.path.join(artifacts_dir, "pptx_salles_seminaire_1785866596542.png")
    img_chambre = os.path.join(artifacts_dir, "pptx_chambre_hotes_1785866613303.png")
    img_dialogue = os.path.join(artifacts_dir, "pptx_intergenerationnel_1785866630658.png")
    img_agri = os.path.join(artifacts_dir, "pptx_agri_innovation_1785866645077.png")
    img_pres = os.path.join(public_images, "professional_corporate_portrait_of_a_west_african_woman_leader_elegant_and_screen.png")
    img_map = os.path.join(public_images, "a_modern_glowing_abstract_3d_map_of_senegal_composed_of_floating_dots_and_nodes_screen.png")

    def apply_transition(slide, trans_type="fade"):
        """Adds native PowerPoint slide transition (fade, push, wipe, zoom, wheel, dissolve)"""
        try:
            xml = f'<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:{trans_type}/></p:transition>'
            slide._element.append(parse_xml(xml))
        except Exception as e:
            print("Transition error:", e)

    def draw_bg(slide, color=SOFT_BG):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = color
        bg.line.fill.background()
        return bg

    def add_light_header(slide, title, category="APATAM@E • VISION 2026", dark_mode=False):
        # Category Tag
        tb = slide.shapes.add_textbox(Inches(0.9), Inches(0.4), Inches(10), Inches(0.35))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = category.upper()
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = BRIGHT_GOLD if dark_mode else OCEAN_BLUE

        # Large Title (Grande Écriture)
        tb2 = slide.shapes.add_textbox(Inches(0.9), Inches(0.75), Inches(10.5), Inches(0.85))
        tf2 = tb2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = title
        p2.font.size = Pt(28)
        p2.font.bold = True
        p2.font.color.rgb = CARD_WHITE if dark_mode else DARK_FOREST

        # Mini Logo top right
        if os.path.exists(logo_path):
            slide.shapes.add_picture(logo_path, Inches(11.8), Inches(0.4), width=Inches(0.75))

    # =========================================================================
    # SLIDE 1: HERO / COVER (Grande Écriture & Animation Fade-Zoom)
    # =========================================================================
    s1 = prs.slides.add_slide(blank_layout)
    draw_bg(s1, DARK_FOREST)
    apply_transition(s1, "fade")

    # Right Hero Image (Clean, Minimalist Framing)
    if os.path.exists(img_hero):
        s1.shapes.add_picture(img_hero, Inches(6.8), Inches(0.8), width=Inches(5.7))

    # Logo Top Left
    if os.path.exists(logo_path):
        s1.shapes.add_picture(logo_path, Inches(0.9), Inches(0.9), width=Inches(1.8))

    # Tag Badge
    tag = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(2.9), Inches(3.6), Inches(0.45))
    tag.fill.solid()
    tag.fill.fore_color.rgb = DEEP_GREEN
    tag.line.color.rgb = ACCENT_GOLD
    tp = tag.text_frame.paragraphs[0]
    tp.text = "PLATEFORME INSTITUTIONNELLE & DÉVELOPPEMENT"
    tp.font.size = Pt(8.5)
    tp.font.bold = True
    tp.font.color.rgb = BRIGHT_GOLD
    tp.alignment = PP_ALIGN.CENTER

    # Huge Title (Grande Écriture 50pt)
    tb_title = s1.shapes.add_textbox(Inches(0.85), Inches(3.5), Inches(5.8), Inches(2.5))
    tf_t = tb_title.text_frame
    tf_t.word_wrap = True

    p_t1 = tf_t.paragraphs[0]
    p_t1.text = "APATAM@E"
    p_t1.font.size = Pt(50)
    p_t1.font.bold = True
    p_t1.font.color.rgb = BRIGHT_GOLD

    p_t2 = tf_t.add_paragraph()
    p_t2.text = "Gouvernance & Territoires d'Avenir"
    p_t2.font.size = Pt(26)
    p_t2.font.bold = True
    p_t2.font.color.rgb = CARD_WHITE

    p_t3 = tf_t.add_paragraph()
    p_t3.text = "\nBâtir des territoires inclusifs, résilients et connectés aux savoirs."
    p_t3.font.size = Pt(13)
    p_t3.font.color.rgb = LIGHT_MINT

    # Footer Metadata
    tb_foot = s1.shapes.add_textbox(Inches(0.9), Inches(6.3), Inches(11.5), Inches(0.6))
    p_f = tb_foot.text_frame.paragraphs[0]
    p_f.text = "Présidente Fondatrice : Adama Mbengue  •  Siège : ZAC NORD Thiès Nº2688  •  contact@apatame.com"
    p_f.font.size = Pt(10)
    p_f.font.color.rgb = CARD_WHITE

    # =========================================================================
    # SLIDE 2: IDENTITÉ & VISION STRATÉGIQUE (Grande Écriture + Transition Push)
    # =========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    draw_bg(s2, SOFT_BG)
    apply_transition(s2, "push")
    add_light_header(s2, "Notre Identité & Notre Raison d'Être", "01 / NOTRE ADN")

    # Left Image with Soft Shadow Border
    if os.path.exists(img_dialogue):
        s2.shapes.add_picture(img_dialogue, Inches(0.9), Inches(1.8), width=Inches(4.6))

    # Right Content: 3 Clean Minimalist Cards with Light Icons
    cards_s2 = [
        ("✦  Mission Fondatrice", "Former, informer, accompagner et mobiliser les citoyens pour une gouvernance responsable et un développement local inclusif et durable."),
        ("✦  Approche Endogène", "Réconcilier les traditions positives africaines avec les innovations technologiques, écologiques et managériales contemporaines."),
        ("✦  Ancrage Territorial", "Un siège national à Thiès (ZAC Nord), des représentations à Dakar et Saly, et des actions dans les 14 régions du Sénégal.")
    ]

    cur_y = 1.8
    for title, desc in cards_s2:
        c = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.8), Inches(cur_y), Inches(6.6), Inches(1.5))
        c.fill.solid()
        c.fill.fore_color.rgb = CARD_WHITE
        c.line.color.rgb = CARD_BORDER
        c.line.width = Pt(1)
        ctf = c.text_frame
        ctf.word_wrap = True
        ctf.margin_left = Inches(0.3)
        ctf.margin_top = Inches(0.2)

        p1 = ctf.paragraphs[0]
        p1.text = title
        p1.font.bold = True
        p1.font.size = Pt(15)
        p1.font.color.rgb = DARK_FOREST

        p2 = ctf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = BODY_TEXT
        cur_y += 1.7

    # =========================================================================
    # SLIDE 3: LE MOT DE LA PRÉSIDENTE (Grand Impact Éditorial)
    # =========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    draw_bg(s3, SOFT_BG)
    apply_transition(s3, "fade")
    add_light_header(s3, "Le Mot de la Présidente Fondatrice", "02 / GOUVERNANCE & LEADERSHIP")

    # President Photo
    if os.path.exists(img_pres):
        s3.shapes.add_picture(img_pres, Inches(0.9), Inches(1.8), width=Inches(3.8))

    # Large Modern Quote Card
    q_card = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.0), Inches(1.8), Inches(7.4), Inches(4.9))
    q_card.fill.solid()
    q_card.fill.fore_color.rgb = DARK_FOREST
    q_card.line.color.rgb = ACCENT_GOLD
    q_card.line.width = Pt(1.5)

    qtf = q_card.text_frame
    qtf.word_wrap = True
    qtf.margin_left = Inches(0.5)
    qtf.margin_right = Inches(0.5)
    qtf.margin_top = Inches(0.5)

    qp1 = qtf.paragraphs[0]
    qp1.text = "« Notre mission est de promouvoir des modèles endogènes de gouvernance et de développement, en réconciliant modernité et traditions positives pour une Afrique debout. »"
    qp1.font.size = Pt(20)
    qp1.font.bold = True
    qp1.font.italic = True
    qp1.font.color.rgb = CARD_WHITE

    qp2 = qtf.add_paragraph()
    qp2.text = "\nÀ travers APATAM@E, nous mobilisons l'intelligence collective des territoires pour que chaque citoyenne et citoyen devienne l'acteur central de son propre épanouissement."
    qp2.font.size = Pt(12.5)
    qp2.font.color.rgb = LIGHT_MINT

    qp3 = qtf.add_paragraph()
    qp3.text = "\n\nAdama Mbengue"
    qp3.font.size = Pt(18)
    qp3.font.bold = True
    qp3.font.color.rgb = BRIGHT_GOLD

    qp4 = qtf.add_paragraph()
    qp4.text = "Présidente Fondatrice — APATAM@E Sénégal"
    qp4.font.size = Pt(11)
    qp4.font.color.rgb = CARD_WHITE

    # =========================================================================
    # SLIDE 4: NOS 7 PRINCIPES D'ACTION & VALEURS (Light Icons & Grilles Épurées)
    # =========================================================================
    s4 = prs.slides.add_slide(blank_layout)
    draw_bg(s4, SOFT_BG)
    apply_transition(s4, "wipe")
    add_light_header(s4, "Nos Principes d'Action & Valeurs Fondatrices", "03 / FONDEMENTS ÉTHIQUES")

    values_data = [
        ("01", "Équité", "Justice sociale et égalité des chances pour tous."),
        ("02", "Intégrité", "Transparence, redevabilité et exemplarité."),
        ("03", "Inclusion", "Participation de tous les terroirs et des femmes."),
        ("04", "Innovation", "Modernité et valorisation des savoirs endogènes."),
        ("05", "Collaboration", "Synergies et partenariats co-construits."),
        ("06", "Paix & Sécurité", "Médiation communautaire et sécurité citoyenne."),
        ("07", "Transmission", "Dialogue intergénérationnel aînés & jeunesse.")
    ]

    # Row 1 (4 items)
    w_card = Inches(2.7)
    h_card = Inches(2.2)
    for i in range(4):
        x = Inches(0.9 + i * 2.95)
        y = Inches(1.8)
        c = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w_card, h_card)
        c.fill.solid()
        c.fill.fore_color.rgb = CARD_WHITE
        c.line.color.rgb = CARD_BORDER
        ctf = c.text_frame
        ctf.word_wrap = True
        ctf.margin_top = Inches(0.2)
        ctf.margin_left = Inches(0.25)
        ctf.margin_right = Inches(0.2)

        p1 = ctf.paragraphs[0]
        p1.text = values_data[i][0]
        p1.font.bold = True
        p1.font.size = Pt(22)
        p1.font.color.rgb = OCEAN_BLUE

        p2 = ctf.add_paragraph()
        p2.text = values_data[i][1]
        p2.font.bold = True
        p2.font.size = Pt(14)
        p2.font.color.rgb = DARK_FOREST

        p3 = ctf.add_paragraph()
        p3.text = values_data[i][2]
        p3.font.size = Pt(10)
        p3.font.color.rgb = BODY_TEXT

    # Row 2 (3 items)
    for j in range(3):
        idx = 4 + j
        x = Inches(1.5 + j * 3.65)
        y = Inches(4.3)
        c = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.35), h_card)
        c.fill.solid()
        c.fill.fore_color.rgb = CARD_WHITE
        c.line.color.rgb = CARD_BORDER
        ctf = c.text_frame
        ctf.word_wrap = True
        ctf.margin_top = Inches(0.2)
        ctf.margin_left = Inches(0.25)
        ctf.margin_right = Inches(0.2)

        p1 = ctf.paragraphs[0]
        p1.text = values_data[idx][0]
        p1.font.bold = True
        p1.font.size = Pt(22)
        p1.font.color.rgb = BRIGHT_GOLD if j == 1 else OCEAN_BLUE

        p2 = ctf.add_paragraph()
        p2.text = values_data[idx][1]
        p2.font.bold = True
        p2.font.size = Pt(14)
        p2.font.color.rgb = DARK_FOREST

        p3 = ctf.add_paragraph()
        p3.text = values_data[idx][2]
        p3.font.size = Pt(10)
        p3.font.color.rgb = BODY_TEXT

    # =========================================================================
    # SLIDE 5: NOS 6 PROGRAMMES STRATÉGIQUES (Design Moderne & Grandes Écritures)
    # =========================================================================
    s5 = prs.slides.add_slide(blank_layout)
    draw_bg(s5, SOFT_BG)
    apply_transition(s5, "push")
    add_light_header(s5, "Nos 6 Programmes & Axes d'Intervention", "04 / PROGRAMMES & IMPACT")

    prog_list = [
        ("01", "Éducation & STEM", "Bourses d'excellence, mentorat scientifique et laboratoires numériques mobiles."),
        ("02", "Leadership Jeunesse", "Incubateur de leadership territorial et renforcement des capacités citoyennes."),
        ("03", "Résilience Agricole", "Agroécologie sahélienne, irrigation solaire et valorisation des filières locales."),
        ("04", "Observatoire Public", "Recherche-action, données de terrain et plaidoyer pour des politiques justes."),
        ("05", "Paix & Sécurité", "Comités locaux de médiation, prévention des conflits et sécurité citoyenne."),
        ("06", "Dialogue Aînés & Jeunes", "Sessions immersives sous l'Apatam et transmission des savoirs endogènes.")
    ]

    pw = Inches(3.7)
    ph = Inches(2.2)
    for idx, (pnum, ptitle, pdesc) in enumerate(prog_list):
        row = idx // 3
        col = idx % 3
        x = Inches(0.9 + col * 4.0)
        y = Inches(1.8 + row * 2.5)

        card = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, pw, ph)
        card.fill.solid()
        if idx in [0, 4]:
            card.fill.fore_color.rgb = DARK_FOREST
            card.line.color.rgb = DARK_FOREST
            num_col = BRIGHT_GOLD
            title_col = CARD_WHITE
            desc_col = LIGHT_MINT
        else:
            card.fill.fore_color.rgb = CARD_WHITE
            card.line.color.rgb = CARD_BORDER
            num_col = OCEAN_BLUE
            title_col = DARK_FOREST
            desc_col = BODY_TEXT

        ctf = card.text_frame
        ctf.word_wrap = True
        ctf.margin_top = Inches(0.2)
        ctf.margin_left = Inches(0.25)
        ctf.margin_right = Inches(0.25)

        p1 = ctf.paragraphs[0]
        p1.text = pnum
        p1.font.bold = True
        p1.font.size = Pt(20)
        p1.font.color.rgb = num_col

        p2 = ctf.add_paragraph()
        p2.text = ptitle
        p2.font.bold = True
        p2.font.size = Pt(14)
        p2.font.color.rgb = title_col

        p3 = ctf.add_paragraph()
        p3.text = pdesc
        p3.font.size = Pt(10)
        p3.font.color.rgb = desc_col

    # =========================================================================
    # SLIDE 6: GOUVERNANCE & ORGANIGRAMME
    # =========================================================================
    s6 = prs.slides.add_slide(blank_layout)
    draw_bg(s6, SOFT_BG)
    apply_transition(s6, "fade")
    add_light_header(s6, "Structure de Gouvernance & Organigramme", "05 / ORGANISATION")

    # Level 1: Conseil d'Orientation Stratégique
    c_cos = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), Inches(1.6), Inches(4.3), Inches(0.85))
    c_cos.fill.solid()
    c_cos.fill.fore_color.rgb = OCEAN_BLUE
    c_cos.line.color.rgb = OCEAN_BLUE
    p = c_cos.text_frame.paragraphs[0]
    p.text = "Conseil d'Orientation Stratégique"
    p.font.bold = True
    p.font.size = Pt(12)
    p.font.color.rgb = CARD_WHITE
    p.alignment = PP_ALIGN.CENTER

    # Level 2: Direction Générale
    c_dg = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), Inches(2.7), Inches(4.3), Inches(0.85))
    c_dg.fill.solid()
    c_dg.fill.fore_color.rgb = DARK_FOREST
    c_dg.line.color.rgb = ACCENT_GOLD
    c_dg.line.width = Pt(1.5)
    p = c_dg.text_frame.paragraphs[0]
    p.text = "Direction Générale (Adama Mbengue)"
    p.font.bold = True
    p.font.size = Pt(12)
    p.font.color.rgb = BRIGHT_GOLD
    p.alignment = PP_ALIGN.CENTER

    c_sec = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.1), Inches(2.7), Inches(3.3), Inches(0.85))
    c_sec.fill.solid()
    c_sec.fill.fore_color.rgb = CARD_WHITE
    c_sec.line.color.rgb = CARD_BORDER
    p = c_sec.text_frame.paragraphs[0]
    p.text = "Secrétariat de Direction\n(Agenda, courrier & accueil)"
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = DARK_FOREST
    p.alignment = PP_ALIGN.CENTER

    # Level 3: 4 Directions Opérationnelles
    dirs_titles = [
        "Dir. des Programmes\n& Formations",
        "Dir. de la Recherche\n& Observatoire",
        "Dir. des Partenariats\n& Communication",
        "Dir. Administrative\n& Financière"
    ]
    for i, d in enumerate(dirs_titles):
        cd = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9 + i * 2.95), Inches(3.8), Inches(2.7), Inches(1.1))
        cd.fill.solid()
        cd.fill.fore_color.rgb = DEEP_GREEN
        cd.line.color.rgb = DEEP_GREEN
        p = cd.text_frame.paragraphs[0]
        p.text = d
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = CARD_WHITE
        p.alignment = PP_ALIGN.CENTER

    # Level 4: 5 Services d'Appui
    supp_list = ["Comptabilité & Finances", "Ressources Humaines", "Informatique & Digital", "Logistique & Patrimoine", "Suivi-Éval & Qualité"]
    for j, s in enumerate(supp_list):
        cs = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9 + j * 2.35), Inches(5.2), Inches(2.2), Inches(0.8))
        cs.fill.solid()
        cs.fill.fore_color.rgb = CARD_WHITE
        cs.line.color.rgb = CARD_BORDER
        p = cs.text_frame.paragraphs[0]
        p.text = s
        p.font.bold = True
        p.font.size = Pt(9)
        p.font.color.rgb = DARK_FOREST
        p.alignment = PP_ALIGN.CENTER

    # Bottom Banner
    b_bar = s6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(6.3), Inches(11.53), Inches(0.6))
    b_bar.fill.solid()
    b_bar.fill.fore_color.rgb = DARK_FOREST
    b_bar.line.fill.background()
    p = b_bar.text_frame.paragraphs[0]
    p.text = "NOTRE MISSION : Former, informer, accompagner et mobiliser pour un développement territorial durable."
    p.font.size = Pt(10.5)
    p.font.bold = True
    p.font.color.rgb = CARD_WHITE
    p.alignment = PP_ALIGN.CENTER

    # =========================================================================
    # SLIDE 7: SALLES DE SÉMINAIRES & CHAMBRES D'HÔTES (Photos Générées Récentes)
    # =========================================================================
    s7 = prs.slides.add_slide(blank_layout)
    draw_bg(s7, SOFT_BG)
    apply_transition(s7, "push")
    add_light_header(s7, "Espaces Mutualisés & Chambres d'Hôtes", "06 / INFRASTRUCTURES & ACCUEIL")

    # Left Column: Salles
    if os.path.exists(img_seminaire):
        s7.shapes.add_picture(img_seminaire, Inches(0.9), Inches(1.8), width=Inches(5.5))
    
    card_salle = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(4.5), Inches(5.5), Inches(2.4))
    card_salle.fill.solid()
    card_salle.fill.fore_color.rgb = CARD_WHITE
    card_salle.line.color.rgb = DARK_FOREST
    card_salle.line.width = Pt(1.5)
    stf = card_salle.text_frame
    stf.word_wrap = True
    stf.margin_left = Inches(0.3)
    stf.margin_top = Inches(0.2)
    sp1 = stf.paragraphs[0]
    sp1.text = "🏛️ Salles de Séminaires & Réunions"
    sp1.font.bold = True
    sp1.font.size = Pt(14)
    sp1.font.color.rgb = DARK_FOREST
    sp2 = stf.add_paragraph()
    sp2.text = "• Siège Thiès ZAC Nord : modulable selon l'effectif exact.\n• Équipements : Vidéoprojecteur, sonorisation, Wifi et climatisation.\n• Activités : Séminaires, Réunions, Assemblées Générales, Formations."
    sp2.font.size = Pt(10)
    sp2.font.color.rgb = BODY_TEXT

    # Right Column: Chambres
    if os.path.exists(img_chambre):
        s7.shapes.add_picture(img_chambre, Inches(6.9), Inches(1.8), width=Inches(5.5))

    card_chambre = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(4.5), Inches(5.5), Inches(2.4))
    card_chambre.fill.solid()
    card_chambre.fill.fore_color.rgb = DARK_FOREST
    card_chambre.line.color.rgb = ACCENT_GOLD
    card_chambre.line.width = Pt(1.5)
    htf = card_chambre.text_frame
    htf.word_wrap = True
    htf.margin_left = Inches(0.3)
    htf.margin_top = Inches(0.2)
    hp1 = htf.paragraphs[0]
    hp1.text = "🏡 Chambres d'Hôtes (Thiès & Saly)"
    hp1.font.bold = True
    hp1.font.size = Pt(14)
    hp1.font.color.rgb = BRIGHT_GOLD
    hp2 = htf.add_paragraph()
    hp2.text = "• Résidence Thiès : calme & sérénité pour missions professionnelles.\n• Villa Saly Portudal : cadre balnéaire, piscine & retraites de travail.\n• Services : Petit-déjeuner bio local inclus, navette Aéroport AIBD."
    hp2.font.size = Pt(10)
    hp2.font.color.rgb = CARD_WHITE

    # =========================================================================
    # SLIDE 8: IMPACT & RÉSULTATS CLÉS (Grandes Écritures Chiffrées 40pt)
    # =========================================================================
    s8 = prs.slides.add_slide(blank_layout)
    draw_bg(s8, DARK_FOREST)
    apply_transition(s8, "fade")
    add_light_header(s8, "Indicateurs de Performance & Impact Réel", "07 / IMPACT & MESURE", dark_mode=True)

    metrics = [
        ("25 000+", "Bénéficiaires Directs", "Agricultrices, lycéennes, jeunes leaders et acteurs territoriaux mobilisés."),
        ("48", "Périmètres Décarbonés", "Systèmes d'irrigation solaire et fermes agroécologiques autonomes."),
        ("94%", "Taux d'Autonomie", "Des coopératives féminines accompagnées pérennes après 24 mois."),
        ("14", "Publications Open Data", "Notes scientifiques et baromètres de gouvernance locale publiés.")
    ]

    for idx, (m_val, m_lbl, m_desc) in enumerate(metrics):
        x = Inches(0.9 + idx * 2.95)
        y = Inches(2.0)
        mc = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.7), Inches(4.5))
        mc.fill.solid()
        mc.fill.fore_color.rgb = DEEP_GREEN
        mc.line.color.rgb = BRIGHT_GOLD if idx == 0 else RGBColor(40, 85, 60)
        mc.line.width = Pt(1.5)

        mtf = mc.text_frame
        mtf.word_wrap = True
        mtf.margin_top = Inches(0.4)
        mtf.margin_left = Inches(0.2)
        mtf.margin_right = Inches(0.2)

        p1 = mtf.paragraphs[0]
        p1.text = m_val
        p1.font.bold = True
        p1.font.size = Pt(36) # Grande écriture
        p1.font.color.rgb = BRIGHT_GOLD
        p1.alignment = PP_ALIGN.CENTER

        p2 = mtf.add_paragraph()
        p2.text = f"\n{m_lbl}"
        p2.font.bold = True
        p2.font.size = Pt(14)
        p2.font.color.rgb = CARD_WHITE
        p2.alignment = PP_ALIGN.CENTER

        p3 = mtf.add_paragraph()
        p3.text = f"\n\n{m_desc}"
        p3.font.size = Pt(10.5)
        p3.font.color.rgb = LIGHT_MINT
        p3.alignment = PP_ALIGN.CENTER

    # =========================================================================
    # SLIDE 9: ALLIANCES & PARTENARIATS STRATÉGIQUES
    # =========================================================================
    s9 = prs.slides.add_slide(blank_layout)
    draw_bg(s9, SOFT_BG)
    apply_transition(s9, "wipe")
    add_light_header(s9, "Nos Alliances & Partenariats Stratégiques", "08 / COOPÉRATION MULTI-ACTEURS")

    p_categories = [
        ("🏛️  Ministères & Agences Publiques", [
            "Ministère de l'Environnement et de la Transition Écologique",
            "Ministère de l'Enseignement Supérieur et de la Recherche",
            "Agence Nationale des Énergies Renouvelables (ANER)",
            "Collectivités territoriales et Communes partenaires"
        ]),
        ("🌍  Bailleurs & Coopération Multilatérale", [
            "Union Européenne (Programmes Climat & Jeunesse)",
            "PNUD Sénégal (Fonds Environnement Mondial)",
            "Coopération Allemande au Développement (GIZ)",
            "Fondations philanthropiques pour l'Afrique"
        ]),
        ("🎓  Universités & Centres de Recherche", [
            "Université Cheikh Anta Diop (UCAD - Dakar)",
            "Université Gaston Berger (UGB - Saint-Louis)",
            "Institut Sénégalais de Recherches Agricoles (ISRA)",
            "Institut de Technologie Alimentaire (ITA)"
        ])
    ]

    for idx, (cat_title, p_list) in enumerate(p_categories):
        x = Inches(0.9 + idx * 3.95)
        y = Inches(1.8)
        p_card = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.7), Inches(4.8))
        p_card.fill.solid()
        p_card.fill.fore_color.rgb = CARD_WHITE
        p_card.line.color.rgb = CARD_BORDER
        
        ptf = p_card.text_frame
        ptf.word_wrap = True
        ptf.margin_top = Inches(0.3)
        ptf.margin_left = Inches(0.3)
        ptf.margin_right = Inches(0.3)

        p1 = ptf.paragraphs[0]
        p1.text = cat_title
        p1.font.bold = True
        p1.font.size = Pt(14)
        p1.font.color.rgb = DARK_FOREST

        for p_item in p_list:
            p_bullet = ptf.add_paragraph()
            p_bullet.text = f"\n• {p_item}"
            p_bullet.font.size = Pt(11)
            p_bullet.font.color.rgb = BODY_TEXT

    # =========================================================================
    # SLIDE 10: CONTACT OFFICIEL & RÉSERVATIONS (Closing Impactant)
    # =========================================================================
    s10 = prs.slides.add_slide(blank_layout)
    draw_bg(s10, DARK_FOREST)
    apply_transition(s10, "fade")
    add_light_header(s10, "Rejoindre l'Action & Nous Contacter", "09 / CONTACT & RÉSERVATIONS", dark_mode=True)

    main_c = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.8), Inches(10.33), Inches(4.8))
    main_c.fill.solid()
    main_c.fill.fore_color.rgb = DEEP_GREEN
    main_c.line.color.rgb = BRIGHT_GOLD
    main_c.line.width = Pt(1.5)

    mtf = main_c.text_frame
    mtf.word_wrap = True
    mtf.margin_left = Inches(0.6)
    mtf.margin_right = Inches(0.6)
    mtf.margin_top = Inches(0.4)

    p1 = mtf.paragraphs[0]
    p1.text = "APATAM@E — Plateforme Institutionnelle"
    p1.font.bold = True
    p1.font.size = Pt(22)
    p1.font.color.rgb = BRIGHT_GOLD

    p2 = mtf.add_paragraph()
    p2.text = "Présidente Fondatrice : Adama Mbengue\n"
    p2.font.size = Pt(14)
    p2.font.bold = True
    p2.font.color.rgb = CARD_WHITE

    p3 = mtf.add_paragraph()
    p3.text = "📍 Siège National : ZAC NORD Thiès Nº2688, Rond Point Mosquée, Thiès - Sénégal"
    p3.font.size = Pt(12.5)
    p3.font.color.rgb = LIGHT_MINT

    p4 = mtf.add_paragraph()
    p4.text = "📧 Courriels : contact@apatame.com  |  secretariat@apatame.org"
    p4.font.size = Pt(12.5)
    p4.font.color.rgb = CARD_WHITE

    p5 = mtf.add_paragraph()
    p5.text = "📞 Téléphones : Bureau : 33 999 28 52  |  Portable : 77 510 20 38"
    p5.font.size = Pt(12.5)
    p5.font.color.rgb = CARD_WHITE

    p6 = mtf.add_paragraph()
    p6.text = "🏡 Espaces Mutualisés & Chambres d'Hôtes : Thiès & Saly Portudal"
    p6.font.size = Pt(12.5)
    p6.font.color.rgb = BRIGHT_GOLD

    p7 = mtf.add_paragraph()
    p7.text = "\n🌐 Plateforme Web en direct : http://localhost:5173"
    p7.font.size = Pt(13)
    p7.font.bold = True
    p7.font.color.rgb = CARD_WHITE

    # Save to all target paths
    output_filename = "Presentation_APATAM_E_Moderne_Animee_2026.pptx"
    
    # 1. Parent folder
    parent_dir = os.path.abspath(os.path.join(base_dir, ".."))
    parent_path = os.path.join(parent_dir, output_filename)
    prs.save(parent_path)
    print(f"[OK] Saved to parent folder: {parent_path}")

    # 2. Downloads folder
    user_home = os.path.expanduser("~")
    downloads_path = os.path.join(user_home, "Downloads", output_filename)
    prs.save(downloads_path)
    print(f"[OK] Saved to Downloads: {downloads_path}")

    # 3. Project folder
    project_path = os.path.join(base_dir, output_filename)
    prs.save(project_path)
    print(f"[OK] Saved to Project: {project_path}")

if __name__ == "__main__":
    build_modern_animated_deck()
